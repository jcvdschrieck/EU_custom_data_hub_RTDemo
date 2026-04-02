#!/usr/bin/env python3
"""
Generate EU_Custom_DataHub_MessageFlow.pptx

Two horizontal swim lanes:
  - EU Custom Data Hub  (top)  — transaction generation through alarm check & agent queue
  - Ireland Investigation (bottom) — agent verdict, routing, and investigation queue

Eight process steps flowing left-to-right, with a cross-lane arrow when a
suspicious IE transaction is forwarded to the agent / Ireland queue.
"""
from __future__ import annotations
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn as _qn
from lxml import etree

# ── Palette ───────────────────────────────────────────────────────────────────
EU_BLUE      = RGBColor(0x00, 0x33, 0x99)
EU_BLUE_LT   = RGBColor(0xD4, 0xDF, 0xF2)
EU_BLUE_MID  = RGBColor(0x66, 0x88, 0xCC)
EU_YELLOW    = RGBColor(0xFF, 0xED, 0x00)
IE_GREEN     = RGBColor(0x16, 0x9B, 0x62)   # Irish green
IE_GREEN_LT  = RGBColor(0xD0, 0xED, 0xDD)
IE_ORANGE    = RGBColor(0xFF, 0x88, 0x3E)
ALARM_RED    = RGBColor(0xCC, 0x22, 0x22)
ALARM_RED_LT = RGBColor(0xFA, 0xE0, 0xE0)
AGENT_PURPLE = RGBColor(0x6A, 0x3D, 0x9A)
AGENT_PUR_LT = RGBColor(0xE8, 0xDD, 0xF5)
DB_TEAL      = RGBColor(0x00, 0x7A, 0x8A)
DB_TEAL_LT   = RGBColor(0xCC, 0xEA, 0xEE)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY   = RGBColor(0xF0, 0xF2, 0xF5)
MID_GREY     = RGBColor(0xAA, 0xAA, 0xAA)
DARK_GREY    = RGBColor(0x22, 0x22, 0x22)
SIM_BLUE     = RGBColor(0x17, 0x6A, 0xA0)
SIM_BLUE_LT  = RGBColor(0xCE, 0xE5, 0xF2)


# ── Primitives ────────────────────────────────────────────────────────────────

def _rect(slide, x, y, w, h,
          fill=WHITE, border=MID_GREY, border_pt=1.0,
          text="", font_size=10, bold=False,
          text_color=DARK_GREY, align=PP_ALIGN.CENTER,
          v_anchor=MSO_ANCHOR.MIDDLE, wrap=True):
    shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(border_pt)
    if text:
        tf = shp.text_frame
        tf.word_wrap = wrap
        tf.vertical_anchor = v_anchor
        p = tf.paragraphs[0]
        p.alignment = align
        for chunk in text.split("\n"):
            if p.runs:
                p = tf.add_paragraph()
                p.alignment = align
            run = p.add_run()
            run.text = chunk
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = text_color
    return shp


def _txt(slide, text, x, y, w, h,
         font_size=10, bold=False, color=DARK_GREY,
         align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for chunk in text.split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = chunk
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def _line(slide, x1, y1, x2, y2,
          color=DARK_GREY, width_pt=1.5,
          arrow_end=True, dashed=False):
    cx = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    cx.line.color.rgb = color
    cx.line.width = Pt(width_pt)
    ln = cx.line._ln
    if dashed:
        prstDash = etree.SubElement(ln, _qn("a:prstDash"))
        prstDash.set("val", "dash")
    if arrow_end:
        tail = ln.find(_qn("a:tailEnd"))
        if tail is None:
            tail = etree.SubElement(ln, _qn("a:tailEnd"))
        tail.set("type", "triangle")
        tail.set("w", "med")
        tail.set("len", "med")
    return cx


# ── Step card ─────────────────────────────────────────────────────────────────

def _step_card(slide, x, y, w, h,
               step_num, title, body,
               hdr_fill, body_fill=WHITE, hdr_text=WHITE):
    HDR = 0.44
    PAD = 0.10

    # Card outline
    _rect(slide, x, y, w, h, fill=body_fill, border=hdr_fill, border_pt=1.5)
    # Header band
    _rect(slide, x, y, w, HDR, fill=hdr_fill, border=hdr_fill, border_pt=0)

    # Step badge
    BADGE = 0.28
    bx = x + 0.09
    by = y + (HDR - BADGE) / 2
    _rect(slide, bx, by, BADGE, BADGE,
          fill=WHITE, border=hdr_fill, border_pt=0,
          text=str(step_num), font_size=9, bold=True,
          text_color=hdr_fill, align=PP_ALIGN.CENTER,
          v_anchor=MSO_ANCHOR.MIDDLE)

    # Header title
    _txt(slide, title,
         x + BADGE + 0.22, y + 0.06,
         w - BADGE - 0.32, HDR - 0.08,
         font_size=8.5, bold=True, color=hdr_text)

    # Body text
    _txt(slide, body,
         x + PAD, y + HDR + 0.09,
         w - PAD * 2, h - HDR - 0.12,
         font_size=7.5, color=DARK_GREY, align=PP_ALIGN.LEFT)


# ── Small pill badge ──────────────────────────────────────────────────────────

def _pill(slide, x, y, w, h, label, fill, text_color=WHITE):
    _rect(slide, x, y, w, h,
          fill=fill, border=fill, border_pt=0,
          text=label, font_size=7, bold=True,
          text_color=text_color, align=PP_ALIGN.CENTER,
          v_anchor=MSO_ANCHOR.MIDDLE)


# ── DB cylinder ───────────────────────────────────────────────────────────────

def _db_box(slide, x, y, w, h, label, fill=DB_TEAL, fill_lt=DB_TEAL_LT):
    _rect(slide, x, y, w, h,
          fill=fill, border=fill, border_pt=0,
          text=label, font_size=7.5, bold=True,
          text_color=WHITE, align=PP_ALIGN.CENTER,
          v_anchor=MSO_ANCHOR.MIDDLE)
    _rect(slide, x, y, w, h * 0.20,
          fill=fill_lt, border=fill, border_pt=0)


# ── Build ─────────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_GREY

    # ── Title bar ─────────────────────────────────────────────────────────────
    TITLE_H = 0.52
    _rect(slide, 0, 0, 13.33, TITLE_H,
          fill=EU_BLUE, border=EU_BLUE,
          text="EU Custom Data Hub — Message Flow: from Transaction Generation to Ireland Investigation Queue",
          font_size=14, bold=True, text_color=WHITE,
          align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.MIDDLE)

    # EU flag stars (simple yellow dot row)
    _rect(slide, 0, 0, 0.14, TITLE_H, fill=EU_YELLOW, border=EU_YELLOW)

    # ── Lane geometry ─────────────────────────────────────────────────────────
    LANE_TOP  = TITLE_H + 0.06
    LABEL_W   = 1.18
    LANE_H    = (7.5 - LANE_TOP - 0.10) / 2    # ≈ 3.41
    EU_LANE_Y = LANE_TOP
    IE_LANE_Y = LANE_TOP + LANE_H

    # Lane backgrounds
    _rect(slide, 0, EU_LANE_Y, 13.33, LANE_H, fill=EU_BLUE_LT, border=EU_BLUE_LT)
    _rect(slide, 0, IE_LANE_Y, 13.33, LANE_H, fill=IE_GREEN_LT, border=IE_GREEN_LT)

    # Lane separator
    _rect(slide, 0, IE_LANE_Y, 13.33, 0.03, fill=MID_GREY, border=MID_GREY)

    # Lane labels
    _rect(slide, 0, EU_LANE_Y, LABEL_W, LANE_H,
          fill=EU_BLUE, border=EU_BLUE,
          text="EU Custom\nData Hub\n\nEuropean\nCommission\nport 8505",
          font_size=8.5, bold=False, text_color=WHITE,
          align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)

    _rect(slide, 0, IE_LANE_Y, LABEL_W, LANE_H,
          fill=IE_GREEN, border=IE_GREEN,
          text="Ireland\nInvestigation\n\nIrish Revenue\nAgent + Queue",
          font_size=8.5, bold=False, text_color=WHITE,
          align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)

    # ── Card geometry — 8 steps ───────────────────────────────────────────────
    CONTENT_X = LABEL_W + 0.10
    CONTENT_W = 13.33 - CONTENT_X - 0.08
    N         = 8
    GAP       = 0.09
    CARD_W    = (CONTENT_W - (N - 1) * GAP) / N    # ≈ 1.47
    PAD_V     = 0.16
    EU_CARD_Y = EU_LANE_Y + PAD_V
    IE_CARD_Y = IE_LANE_Y + PAD_V
    CARD_H    = LANE_H - 2 * PAD_V                 # ≈ 3.09

    def cx(i):
        return CONTENT_X + i * (CARD_W + GAP)

    EU_MID = EU_LANE_Y + LANE_H / 2
    IE_MID = IE_LANE_Y + LANE_H / 2

    # ── STEP 1 — EU: Simulation DB / Transaction Generation ───────────────────
    _step_card(slide, cx(0), EU_CARD_Y, CARD_W, CARD_H,
               step_num=1,
               title="Simulation DB",
               body=("1,514 March 2026\ntransactions pre-seeded.\n\n"
                     "Scenario week 2\n(8–14 Mar):\n"
                     "TechZone GmbH → IE\n8 txn/day at 0%\n(correct: 23%)\n\n"
                     "Replay speed:\nconfigurab. 1×–2880×"),
               hdr_fill=SIM_BLUE,
               body_fill=SIM_BLUE_LT)

    # simulation.db label
    DB_W = CARD_W - 0.18
    DB_H = 0.44
    _db_box(slide, cx(0) + 0.09, EU_CARD_Y + CARD_H - DB_H - 0.08,
            DB_W, DB_H, "simulation.db", fill=SIM_BLUE, fill_lt=SIM_BLUE_LT)

    # ── STEP 2 — EU: Simulation Engine ───────────────────────────────────────
    _step_card(slide, cx(1), EU_CARD_Y, CARD_W, CARD_H,
               step_num=2,
               title="Simulation Engine",
               body=("asyncio background\ntask; ticks every 50ms.\n\n"
                     "Advances simulated\nclock at configured\nspeed.\n\n"
                     "Fetches pending\ntransactions ≤ sim_time\nfrom simulation.db;\n"
                     "marks fired=1.\n\n"
                     "Calls fire_callback()"),
               hdr_fill=SIM_BLUE,
               body_fill=SIM_BLUE_LT)

    # ── STEP 3 — EU: Insert to European Custom DB ────────────────────────────
    _step_card(slide, cx(2), EU_CARD_Y, CARD_W, CARD_H,
               step_num=3,
               title="Insert to EU Custom DB",
               body=("insert_transaction()\nwrites each record to\n"
                     "european_custom.db.\n\n"
                     "Transaction added to\nin-memory live queue\n"
                     "(30-tx ring buffer\nfor front-end poll).\n\n"
                     "All records stored\nregardless of VAT\ncorrectness."),
               hdr_fill=EU_BLUE,
               body_fill=WHITE)

    DB3_W = CARD_W - 0.18
    DB3_H = 0.44
    _db_box(slide, cx(2) + 0.09, EU_CARD_Y + CARD_H - DB3_H - 0.08,
            DB3_W, DB3_H, "european_custom.db")

    # ── STEP 4 — EU: VAT Alarm Checker ───────────────────────────────────────
    _step_card(slide, cx(3), EU_CARD_Y, CARD_W, CARD_H,
               step_num=4,
               title="VAT Alarm Checker",
               body=("After DB write:\n\n"
                     "7-day VAT/value ratio\nvs 8-week baseline.\n\n"
                     "Deviation > 25%:\n→ alarm raised\n   (7-day validity)\n\n"
                     "Active alarm + IE:\n→ suspicious = 1\n   suspicion: MEDIUM\n\n"
                     "Non-IE: alarm only,\nno suspicious flag."),
               hdr_fill=ALARM_RED,
               body_fill=ALARM_RED_LT)

    # alarm badge
    _pill(slide,
          cx(3) + 0.09, EU_CARD_Y + CARD_H - 0.34,
          CARD_W - 0.18, 0.26,
          "⚠ alarms table", ALARM_RED)

    # ── STEP 5 — EU: Agent Queue ──────────────────────────────────────────────
    _step_card(slide, cx(4), EU_CARD_Y, CARD_W, CARD_H,
               step_num=5,
               title="Agent Queue",
               body=("If suspicious=1\n(IE-bound only):\ntx enqueued to\n"
                     "asyncio.Queue\n(non-blocking).\n\n"
                     "Background worker\npicks up items;\nruns agent in\n"
                     "ThreadPoolExecutor\n(keeps sim loop free).\n\n"
                     "Queue depth shown\nin /api/simulation/\nstatus."),
               hdr_fill=AGENT_PURPLE,
               body_fill=AGENT_PUR_LT)

    # ── STEP 6 — IE: VAT Fraud Detection Agent ───────────────────────────────
    _step_card(slide, cx(5), IE_CARD_Y, CARD_W, CARD_H,
               step_num=6,
               title="VAT Fraud Detection Agent",
               body=("Subprocess bridge:\nagent_bridge.py calls\n"
                     "_analyse_tx.py in\nvat_fraud_detection/\n\n"
                     "Builds Invoice +\nLineItem from tx dict.\n\n"
                     "RAG retrieves Irish\nVAT legislation from\nChromaDB.\n\n"
                     "LM Studio LLM returns\nverdict per line item:\ncorrect / incorrect\n/ uncertain"),
               hdr_fill=AGENT_PURPLE,
               body_fill=AGENT_PUR_LT)

    # LM Studio badge
    _pill(slide,
          cx(5) + 0.09, IE_CARD_Y + CARD_H - 0.34,
          CARD_W - 0.18, 0.26,
          "LM Studio  port 1234", AGENT_PURPLE)

    # ── STEP 7 — IE: Verdict Routing ─────────────────────────────────────────
    _step_card(slide, cx(6), IE_CARD_Y, CARD_W, CARD_H,
               step_num=7,
               title="Verdict Routing",
               body=("INCORRECT:\n"
                     "→ suspicion_level\n   upgraded to HIGH\n"
                     "→ insert_ireland_\n   queue()\n"
                     "→ insert_agent_\n   log(sent=1)\n\n"
                     "CORRECT /\nUNCERTAIN:\n"
                     "→ suspicious flag\n   cleared\n"
                     "→ insert_agent_\n   log(sent=0)"),
               hdr_fill=IE_GREEN,
               body_fill=IE_GREEN_LT)

    # Verdict pills
    _pill(slide, cx(6) + 0.09, IE_CARD_Y + CARD_H - 0.62,
          CARD_W - 0.18, 0.24, "INCORRECT → HIGH", ALARM_RED)
    _pill(slide, cx(6) + 0.09, IE_CARD_Y + CARD_H - 0.34,
          CARD_W - 0.18, 0.24, "CORRECT / UNCERTAIN → cleared",
          IE_GREEN)

    # ── STEP 8 — IE: Ireland Investigation Queue ──────────────────────────────
    _step_card(slide, cx(7), IE_CARD_Y, CARD_W, CARD_H,
               step_num=8,
               title="Ireland Investigation Queue",
               body=("Confirmed-incorrect\ncases available for\nlocal investigation.\n\n"
                     "Each entry includes:\n· Transaction detail\n"
                     "· Alarm key + deviation\n"
                     "· Applied vs correct\n  VAT rate\n"
                     "· Agent reasoning\n· Suspicion: HIGH\n\n"
                     "Front-end:\n/ireland page"),
               hdr_fill=IE_GREEN,
               body_fill=IE_GREEN_LT)

    DB8_W = CARD_W - 0.18
    DB8_H = 0.44
    _db_box(slide, cx(7) + 0.09, IE_CARD_Y + CARD_H - DB8_H - 0.08,
            DB8_W, DB8_H, "ireland_queue table",
            fill=IE_GREEN, fill_lt=IE_GREEN_LT)

    # ── Within-EU lane arrows (steps 1→2→3→4→5) ──────────────────────────────
    for i in range(4):
        _line(slide, cx(i) + CARD_W, EU_MID, cx(i + 1), EU_MID,
              color=EU_BLUE, width_pt=2.0)

    # ── Within-Ireland lane arrows (steps 6→7→8) ─────────────────────────────
    for i in [5, 6]:
        _line(slide, cx(i) + CARD_W, IE_MID, cx(i + 1), IE_MID,
              color=IE_GREEN, width_pt=2.0)

    # ── Cross-lane arrow: Agent Queue → Agent (suspicious IE tx) ─────────────
    LANE_SEP    = IE_LANE_Y
    EU_CARD_BTM = EU_CARD_Y + CARD_H
    IE_CARD_TOP = IE_CARD_Y

    # Drop from step 5 (Agent Queue, EU lane) to step 6 (Agent, IE lane)
    # Horizontal mid of step 5 → horizontal mid of step 6
    CROSS_X5 = cx(4) + CARD_W * 0.55   # step 5 drop point (EU lane)
    CROSS_X6 = cx(5) + CARD_W * 0.45   # step 6 receive point (IE lane)

    # Elbow: down from step 5, across to step 6, then down to step 6 top
    # Segment 1: step 5 bottom → lane separator
    _line(slide, CROSS_X5, EU_CARD_BTM, CROSS_X5, LANE_SEP,
          color=AGENT_PURPLE, width_pt=2.0, arrow_end=False)
    # Segment 2: horizontal across lane separator
    _line(slide, CROSS_X5, LANE_SEP, CROSS_X6, LANE_SEP,
          color=AGENT_PURPLE, width_pt=2.0, arrow_end=False)
    # Segment 3: down to step 6 top
    _line(slide, CROSS_X6, LANE_SEP, CROSS_X6, IE_CARD_TOP,
          color=AGENT_PURPLE, width_pt=2.0, arrow_end=True)

    # Label on the cross-lane path (in the gap between lanes)
    GAP_LBL_Y = LANE_SEP - 0.38
    _txt(slide,
         "suspicious IE tx\n(suspicion: MEDIUM)",
         min(CROSS_X5, CROSS_X6) - 0.05, GAP_LBL_Y, 1.30, 0.38,
         font_size=7, color=AGENT_PURPLE, italic=True,
         align=PP_ALIGN.CENTER)

    # ── Also note: EU Custom DB stores everything (arrow from step 3 label) ────
    # "All txn stored" note box between step 3 (EU) label
    NOTE_X = cx(2) + CARD_W * 0.10
    NOTE_Y = EU_CARD_Y - 0.01
    # (already described in card body — no extra note needed)

    # ── European Custom DB note box (spanning steps 3–5) ─────────────────────
    # Show that agent_log also lands in the same DB
    NOTE_DB_X = cx(2) + 0.04
    NOTE_DB_Y = EU_CARD_Y + CARD_H + 0.03   # just below the EU cards
    NOTE_DB_W = cx(4) + CARD_W - NOTE_DB_X  # spans steps 3 through 5
    NOTE_DB_H = 0.26

    # Only show if there is space (there's very little room in 2-lane layout)
    # — skip; described in card bodies already —

    # ── Suspicion-level legend ────────────────────────────────────────────────
    # Small pills in the gap between lanes at the right
    LEG_X = cx(7) + CARD_W + 0.04
    if LEG_X < 13.20:
        pass   # no room; legend goes in footer

    # ── Footer legend ─────────────────────────────────────────────────────────
    LEG_Y = 7.5 - 0.26
    _txt(slide,
         ("─── EU Data Hub flow (blue)     "
          "─── Ireland investigation flow (green)     "
          "─── Agent handoff: suspicious IE transaction (purple)     "
          "All transactions recorded in european_custom.db regardless of flow"),
         0.18, LEG_Y, 13.0, 0.24,
         font_size=6.5, color=MID_GREY, align=PP_ALIGN.LEFT)

    # ── Suspicion level note in the IE label ─────────────────────────────────
    # Add small note under the IE lane label
    NOTE_W = LABEL_W - 0.10
    NOTE_H = 0.46
    _rect(slide,
          0.05, IE_LANE_Y + LANE_H - NOTE_H - 0.05, NOTE_W, NOTE_H,
          fill=ALARM_RED_LT, border=ALARM_RED, border_pt=0.8,
          text="MEDIUM\n→ alarm\nHIGH\n→ agent",
          font_size=6.5, bold=False, text_color=ALARM_RED,
          align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)

    # ── Save ──────────────────────────────────────────────────────────────────
    out = "EU_Custom_DataHub_MessageFlow.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    build()
