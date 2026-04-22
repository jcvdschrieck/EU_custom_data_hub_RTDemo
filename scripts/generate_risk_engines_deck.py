#!/usr/bin/env python3
"""
Generate a PowerPoint deck documenting the four risk monitoring engines
plus the consolidation / contribution / case-averaging mechanics.

Output: docs/risk_engines.pptx

Structure:
  1. Title
  2. Overview (four engines + broker + release factory)
  3. Engine 1 — VAT Rate Mismatch (vat_ratio)            + shortcuts box
  4. Engine 2 — Supplier / ML Risk (watchlist)           + shortcuts box
  5. Engine 3 — Ireland Watchlist (ireland_watchlist)    + shortcuts box
  6. Engine 4 — Description Vagueness                    + shortcuts box
  7. Score consolidation (weighted sum + vat_ratio floor + cap)
  8. Routing thresholds (release / investigate / retain)
  9. Contribution share & case-level averaging (non-linearity caveat)
"""
from pathlib import Path

from pptx                import Presentation
from pptx.dml.color      import RGBColor
from pptx.enum.shapes    import MSO_SHAPE
from pptx.enum.text      import PP_ALIGN, MSO_ANCHOR
from pptx.util           import Inches, Pt

# ── Palette ────────────────────────────────────────────────────────────────
EU_BLUE       = RGBColor(0x00, 0x33, 0x99)
EU_GOLD       = RGBColor(0xFF, 0xCC, 0x00)
DARK_TEXT     = RGBColor(0x22, 0x22, 0x22)
LIGHT_TEXT    = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_GREEN  = RGBColor(0x1F, 0x7A, 0x3C)
ACCENT_AMBER  = RGBColor(0xE6, 0x82, 0x0A)
ACCENT_RED    = RGBColor(0xC0, 0x39, 0x2B)
SUBTLE_GREY   = RGBColor(0x66, 0x66, 0x66)
BG_LIGHT      = RGBColor(0xF5, 0xF5, 0xF5)
SHORTCUT_BG   = RGBColor(0xFF, 0xF3, 0xE0)        # light amber
SHORTCUT_EDGE = RGBColor(0xE6, 0x82, 0x0A)        # amber
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Drawing helpers ────────────────────────────────────────────────────────

def _add_text_box(slide, left, top, width, height, text, *,
                  font_size=14, bold=False, color=DARK_TEXT,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                  font_name="Calibri", fill=None, edge=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if edge is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = edge
        shape.line.width = Pt(1.25)
    tf = shape.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return shape


def _add_bullet_box(slide, left, top, width, height, title, bullets, *,
                    title_size=16, body_size=12, fill=None, edge=None,
                    title_color=EU_BLUE, body_color=DARK_TEXT,
                    title_bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if edge is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = edge
        shape.line.width = Pt(1.25)
    tf = shape.text_frame
    tf.margin_left   = Inches(0.15)
    tf.margin_right  = Inches(0.15)
    tf.margin_top    = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    tf.word_wrap = True

    # Title
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    run = p0.add_run()
    run.text = title
    run.font.bold = title_bold
    run.font.size = Pt(title_size)
    run.font.color.rgb = title_color
    run.font.name = "Calibri"

    # Bullets
    for b in bullets:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        run = p.add_run()
        run.text = "• " + b
        run.font.size = Pt(body_size)
        run.font.color.rgb = body_color
        run.font.name = "Calibri"


def _add_rect(slide, left, top, width, height, *, fill, edge=None,
              text=None, font_size=14, bold=False, color=DARK_TEXT,
              align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if edge is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = edge
        shape.line.width = Pt(1)
    if text is not None:
        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return shape


def _add_arrow(slide, x1, y1, x2, y2, color=SUBTLE_GREY):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line = connector.line
    line.color.rgb = color
    line.width = Pt(1.5)
    return connector


def _slide_header(slide, title, subtitle=None):
    _add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.55),
                  title, font_size=28, bold=True, color=EU_BLUE)
    if subtitle:
        _add_text_box(slide, Inches(0.5), Inches(0.75), Inches(12.3), Inches(0.4),
                      subtitle, font_size=14, color=SUBTLE_GREY)


# ── Slide builders ─────────────────────────────────────────────────────────

def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=EU_BLUE)
    _add_text_box(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.0),
                  "Risk Monitoring Engines",
                  font_size=52, bold=True, color=LIGHT_TEXT, align=PP_ALIGN.LEFT)
    _add_text_box(s, Inches(0.8), Inches(3.5), Inches(11.7), Inches(0.6),
                  "EU Custom Data Hub — how transactions are scored",
                  font_size=22, color=EU_GOLD, align=PP_ALIGN.LEFT)
    _add_text_box(s, Inches(0.8), Inches(4.3), Inches(11.7), Inches(1.3),
                  "Four independent engines feed a weighted consolidator that decides "
                  "release / investigate / retain for every incoming sales order.",
                  font_size=16, color=LIGHT_TEXT, align=PP_ALIGN.LEFT)
    _add_text_box(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.4),
                  "Four engines · weighted sum · vat_ratio floor · case-level running mean",
                  font_size=12, color=EU_GOLD, align=PP_ALIGN.LEFT)


def slide_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_header(s, "Overview — four engines, one release factory",
                  "All four subscribe to SALES_ORDER_EVENT and publish to RT_RISK_OUTCOME; "
                  "the Release Factory consolidates and routes.")

    # Source node
    _add_rect(s, Inches(0.6), Inches(3.0), Inches(2.0), Inches(1.0),
              fill=BG_LIGHT, edge=SUBTLE_GREY,
              text="SALES_ORDER_EVENT\nbroker",
              font_size=11, bold=True)

    # Four engine boxes
    engine_names = [
        ("vat_ratio",           "VAT Rate\nMismatch"),
        ("watchlist",           "Supplier /\nML Risk"),
        ("ireland_watchlist",   "Ireland\nWatchlist"),
        ("description_vagueness","Description\nVagueness"),
    ]
    eng_x = Inches(3.8)
    eng_y0 = Inches(1.2)
    eng_h = Inches(1.0)
    eng_w = Inches(2.5)
    gap = Inches(0.25)
    for i, (eid, label) in enumerate(engine_names):
        _add_rect(s, eng_x, eng_y0 + i * (eng_h + gap), eng_w, eng_h,
                  fill=EU_BLUE, edge=None, text=label,
                  font_size=13, bold=True, color=LIGHT_TEXT)
        _add_text_box(s,
                      Inches(3.8 + 2.55), eng_y0 + i * (eng_h.emu + gap.emu) / 914400 * 914400,
                      Inches(2.5), Inches(0.3),
                      f"engine = \"{eid}\"",
                      font_size=9, color=SUBTLE_GREY)

    # Consolidator
    cons_x = Inches(9.4)
    _add_rect(s, cons_x, Inches(2.9), Inches(3.2), Inches(1.6),
              fill=EU_GOLD, edge=EU_BLUE,
              text="Release Factory\n_compute_score\n(weighted sum ≤ 1)",
              font_size=13, bold=True)

    # Route outputs
    r_x = Inches(9.4)
    for i, (label, color) in enumerate([
        ("release  (< 33.3 %)",    ACCENT_GREEN),
        ("investigate  (33.3–80 %)", ACCENT_AMBER),
        ("retain  (≥ 80 %)",       ACCENT_RED),
    ]):
        _add_rect(s, r_x, Inches(4.8) + i * Inches(0.55),
                  Inches(3.2), Inches(0.5),
                  fill=color, text=label,
                  font_size=12, bold=True, color=LIGHT_TEXT)

    # Arrows broker → engines
    for i in range(4):
        _add_arrow(s,
                   Inches(2.6), Inches(3.5),
                   eng_x, eng_y0 + i * (eng_h + gap) + Inches(0.5))

    # Arrows engines → consolidator
    for i in range(4):
        _add_arrow(s,
                   eng_x + eng_w, eng_y0 + i * (eng_h + gap) + Inches(0.5),
                   cons_x, Inches(3.7))

    # Arrow consolidator → routes
    _add_arrow(s, cons_x + Inches(1.6), Inches(4.5),
               cons_x + Inches(1.6), Inches(4.8))


def engine_slide(prs, *, number, name, engine_id, detects, algorithm,
                 output, weight, shortcuts):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_header(s, f"Engine {number} — {name}",
                  f"engine_id: \"{engine_id}\"   ·   ENGINE_WEIGHTS[{engine_id}] = {weight}")

    # Left column: what / algo / output
    _add_bullet_box(s, Inches(0.5), Inches(1.4), Inches(6.3), Inches(1.5),
                    "What it detects", detects, fill=WHITE, edge=EU_BLUE)
    _add_bullet_box(s, Inches(0.5), Inches(3.0), Inches(6.3), Inches(2.3),
                    "Algorithm (production path)", algorithm, fill=WHITE, edge=EU_BLUE)
    _add_bullet_box(s, Inches(0.5), Inches(5.4), Inches(6.3), Inches(1.6),
                    "Output on RT_RISK_OUTCOME", output, fill=WHITE, edge=EU_BLUE)

    # Right column: SHORTCUTS / DEMO FALLBACKS
    _add_bullet_box(s, Inches(7.0), Inches(1.4), Inches(5.8), Inches(5.6),
                    "Demo shortcuts & fallbacks", shortcuts,
                    fill=SHORTCUT_BG, edge=SHORTCUT_EDGE,
                    title_color=SHORTCUT_EDGE, title_size=16, body_size=11)


def slide_engine_vat_ratio(prs):
    engine_slide(
        prs,
        number=1,
        name="VAT Rate Mismatch",
        engine_id="vat_ratio",
        weight=0.5,
        detects=[
            "Whether the declared VAT rate on an invoice matches the rate the EU expects "
            "for the declared product subcategory in the destination country.",
            "The canonical VAT-misclassification fraud pattern: declaring a low-rate "
            "subcategory (e.g. 'Basic food 0 %') when the product is actually standard-rated.",
        ],
        algorithm=[
            "Key the tx by (buyer_country, vat_subcategory_code).",
            "Look up the expected rate via lib.vat_dataset.expected_rate_for(...). "
            "The lookup is built from the 78 observed (dest, subcat) pairs in the xlsx; "
            "for unobserved pairs it falls back to the destination's standard rate.",
            "Compare declared_rate to expected_rate — emit 1.0 on mismatch, 0.0 on match.",
            "Floor behaviour at consolidation: when raw risk ≥ 0.30 the engine's "
            "contribution is floored to ≥ 0.334 so any genuine rate mismatch at least "
            "investigates, regardless of weighting.",
        ],
        output=[
            "risk ∈ [0, 1] — graded (raw = Score 1 / 100 from xlsx) when pre-baked; "
            "binary 0/1 from the subcategory lookup otherwise.",
            "reason: prebaked / rate_mismatch / rate_match / unknown_subcategory / alarm_*",
            "applicable: always True.",
        ],
        shortcuts=[
            "Pre-baked value: every tx in simulation.db carries "
            "engine_vat_ratio_risk set at seed time from xlsx Score 1 ÷ 100, "
            "GATED by actual rate mismatch (xlsx zeroes Score 1 when declared = "
            "recommended, i.e. misclassification with no revenue impact).",
            "Legacy fallback (historical seeder only): 7-day vs 8-week VAT/value "
            "volume-ratio alarm. Kept so Sep–Feb historical transactions still "
            "route through the pipeline without pre-baked fields.",
            "Demo benefit: routing is deterministic and reproducible from the xlsx — "
            "no need for the seeder to compute rate tables at runtime.",
        ],
    )


def slide_engine_ml(prs):
    engine_slide(
        prs,
        number=2,
        name="Supplier / ML Risk",
        engine_id="watchlist",
        weight=0.9,
        detects=[
            "Per-tx supplier compliance risk: how often has this (seller, origin, "
            "category, destination) pattern been flagged historically?",
            "Produces a graded score plus four per-dimension contributors "
            "(seller, origin, category, destination) that surface in the case.",
        ],
        algorithm=[
            "Look up the 4-tuple (seller, origin, declared_category, destination) "
            "against lib.database.ml_risk_rules.",
            "If a rule matches: emit rule.risk plus the four weighted contributors.",
            "If no rule matches: emit risk = 0.0 (clear).",
            "Flag threshold: ML_RISK_FLAG_THRESHOLD = 0.5 (informational; the "
            "consolidator uses the raw risk).",
        ],
        output=[
            "risk ∈ [0, 1] — graded, discrete buckets from xlsx (0, 0.40, 0.90).",
            "description, seller_risk, country_risk, product_category_risk, "
            "destination_risk — propagated into ASSESSMENT_OUTCOME and written onto "
            "Sales_Order_Risk at case creation.",
            "reason: prebaked_match / prebaked_clear / ml_watchlist_match / clear",
        ],
        shortcuts=[
            "Pre-baked value: engine_ml_risk set per-tx at seed time from xlsx "
            "Score 3 ÷ 100 (0, 0.40, or 0.90). Contributor weights also pre-baked — "
            "all attributed to the seller dimension (seller_contribution = 1.0) since "
            "xlsx Score 3 captures supplier risk specifically.",
            "Legacy fallback: 4-tuple rule lookup in ml_risk_rules table, seeded from "
            "Context/Fake ML.xlsx. Used only when pre-baked field is NULL (legacy "
            "historical tx).",
            "Demo benefit: avoids training or shipping an actual model; "
            "xlsx Score 3 is the ground-truth 'ML output'.",
        ],
    )


def slide_engine_ie(prs):
    engine_slide(
        prs,
        number=3,
        name="Ireland Watchlist",
        engine_id="ireland_watchlist",
        weight=1.0,
        detects=[
            "Country-specific signal hosted (in production) on a server managed by "
            "the Irish authority.",
            "Checks whether a given (seller_id, seller_country) pair is on "
            "Ireland's local blacklist.",
        ],
        algorithm=[
            "Subscribe to every SALES_ORDER_EVENT but only PROCESS events whose "
            "buyer_country == 'IE'.",
            "Non-IE events are immediately published with applicable = False and "
            "excluded from the consolidator's denominator.",
            "Apply a random 1–5 second sleep simulating the round-trip to a "
            "remote server. Because that latency can exceed ASSESSMENT_TIMER_S (3 s), "
            "some IE outcomes legitimately arrive after the consolidator has already "
            "published — by design.",
            "Look up (seller_id, seller_country) in the IE_WATCHLIST frozenset; "
            "emit 1.0 on match, 0.0 on miss.",
        ],
        output=[
            "risk ∈ [0, 1] — binary (1.0 match / 0.0 miss) in production; "
            "graded 0.0 pre-bake in the current dataset.",
            "applicable = False for non-IE tx (engine drops out of denominator).",
            "reason: prebaked_clear / ie_watchlist_match / clear / not_applicable",
        ],
        shortcuts=[
            "Static set IE_WATCHLIST in lib/watchlist.py is currently EMPTY. All "
            "IE-destined tx land in 'clear'.",
            "Pre-baked value: engine_ie_watchlist_risk = 0.0 for every IE tx at "
            "seed time so the engine stays on the pre-baked path rather than "
            "hitting the empty set.",
            "Demo benefit: demonstrates the country-specific-engine pattern "
            "(latency, applicability, remote hosting) without actually importing an "
            "Irish authority feed. The empty watchlist keeps the engine visible in "
            "every case breakdown at 0 %.",
        ],
    )


def slide_engine_vagueness(prs):
    engine_slide(
        prs,
        number=4,
        name="Description Vagueness",
        engine_id="description_vagueness",
        weight=0.8,
        detects=[
            "How generic or ambiguous the product description on the sales order is.",
            "A high score alone is not conclusive, but combined with VAT or supplier "
            "signals it tips the consolidated score over the investigate threshold.",
        ],
        algorithm=[
            "Embed the product description with sentence-transformers/"
            "all-MiniLM-L6-v2 (pre-loaded on backend startup in a worker thread).",
            "Cosine-similarity against a pre-computed anchor embedding built from "
            "a bag of generic phrases ('general goods', 'miscellaneous items', …).",
            "Clamp the raw similarity to [0, 1] and publish as risk.",
            "Flag threshold: risk ≥ 0.5.",
        ],
        output=[
            "risk ∈ [0, 1] — continuous in production, jittered from a binary xlsx "
            "signal in the demo (see shortcuts).",
            "reason: prebaked_clear / prebaked_vague / clear / vague_description / "
            "missing_description",
            "applicable: always True.",
        ],
        shortcuts=[
            "Pre-baked value: engine_vagueness_risk set per-tx at seed time from "
            "xlsx Score 2 ÷ 100 (binary: 0 or 0.60).",
            "Continuous-baseline jitter in lib/new_seeder._jitter_vagueness: "
            "quiet pre-bake → uniform(0.02, 0.08), firing pre-bake → "
            "pre-baked + uniform(-0.03, 0.05). Keeps the UI from showing hard 0 % "
            "while preserving route distribution (max jitter contribution 0.064 ≪ "
            "threshold gaps).",
            "Within-case diversity: the seeder promotes a configurable fraction of "
            "siblings in each investigate cluster to 'vague variants' — generic "
            "description + high vagueness score, low vat_ratio / ml — so case "
            "averages show mixed signals.",
            "Embedding-model fallback: kept for tx with no pre-baked value (legacy "
            "seeder). Slow path but fully continuous.",
        ],
    )


def slide_consolidation(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_header(s, "Score consolidation — weighted sum ≤ 1",
                  "api.py _release_factory._compute_score")

    formula = (
        "score = min( 1.0 ,  Σ  ENGINE_WEIGHTS[eng] × engine.risk   for every "
        "APPLICABLE engine )"
    )
    _add_rect(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.8),
              fill=BG_LIGHT, edge=EU_BLUE, text=formula,
              font_size=16, bold=True, color=EU_BLUE)

    # Weights table
    weights = [
        ("vat_ratio",             "0.5"),
        ("watchlist  (ML)",        "0.9"),
        ("ireland_watchlist",     "1.0"),
        ("description_vagueness", "0.8"),
    ]
    _add_text_box(s, Inches(0.6), Inches(2.5), Inches(5.0), Inches(0.4),
                  "ENGINE_WEIGHTS", font_size=15, bold=True, color=EU_BLUE)
    for i, (name, w) in enumerate(weights):
        y = Inches(2.9 + i * 0.5)
        _add_rect(s, Inches(0.6), y, Inches(3.4), Inches(0.45),
                  fill=WHITE, edge=SUBTLE_GREY, text=name,
                  font_size=12, align=PP_ALIGN.LEFT)
        _add_rect(s, Inches(4.0), y, Inches(1.0), Inches(0.45),
                  fill=EU_GOLD, edge=SUBTLE_GREY, text=w,
                  font_size=12, bold=True)

    _add_bullet_box(s, Inches(5.8), Inches(2.5), Inches(7.0), Inches(2.7),
                    "vat_ratio floor", [
                        "When a tx's raw vat_ratio ≥ 0.30, the engine's contribution is "
                        "floored to ≥ THRESHOLD_RELEASE + ε (≈ 0.334).",
                        "Policy stance: any genuine rate mismatch above the xlsx's "
                        "release tier (Score 1 ≥ 30) must at least INVESTIGATE, "
                        "regardless of its weighted contribution.",
                        "Trigger (0.30) sits between xlsx Release tier (Score 1 = 25) "
                        "and Investigate tier (Score 1 ≥ 37.5). No row that xlsx "
                        "releases gets unexpectedly bumped up.",
                    ], fill=SHORTCUT_BG, edge=SHORTCUT_EDGE, title_color=SHORTCUT_EDGE)

    _add_bullet_box(s, Inches(0.6), Inches(5.5), Inches(12.2), Inches(1.8),
                    "Applicability & confidence", [
                        "Engines may self-report applicable = False (e.g. IE watchlist on "
                        "non-IE tx). Those engines are excluded from the sum AND from the "
                        "denominator used for confidence.",
                        "confidence = applicable_received / applicable_expected. 0 %, 25 %, "
                        "50 %, 75 %, or 100 %.",
                        "If zero applicable engines report, score defaults to 0.5 (uncertain).",
                    ], fill=WHITE, edge=SUBTLE_GREY, body_size=12)


def slide_thresholds(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_header(s, "Routing thresholds",
                  "THRESHOLD_RELEASE = 1/3   ·   THRESHOLD_RETAIN = 0.80")

    rows = [
        (ACCENT_GREEN, "score < 33.33 %",       "Release",     "Terminal event published; no case opened."),
        (ACCENT_AMBER, "33.33 % ≤ score < 80 %", "Investigate", "Sent to the C&T Risk Management Factory; a case is opened in investigation.db and pushed over SSE to the Customs/Tax frontends."),
        (ACCENT_RED,   "score ≥ 80 %",           "Retain",      "Terminal event published; NO case opened (retain bypasses C&T by design — retentions are the outcome of officer escalation from an existing investigate case)."),
    ]
    for i, (color, cond, label, desc) in enumerate(rows):
        y = Inches(1.4 + i * 1.5)
        _add_rect(s, Inches(0.6), y, Inches(2.8), Inches(1.3),
                  fill=color, text=label,
                  font_size=22, bold=True, color=LIGHT_TEXT)
        _add_rect(s, Inches(3.5), y, Inches(3.0), Inches(1.3),
                  fill=BG_LIGHT, edge=SUBTLE_GREY, text=cond,
                  font_size=16, bold=True)
        _add_text_box(s, Inches(6.7), y + Inches(0.1),
                      Inches(6.1), Inches(1.2),
                      desc, font_size=13, color=DARK_TEXT)

    _add_bullet_box(s, Inches(0.6), Inches(6.2), Inches(12.2), Inches(1.2),
                    "Threshold rationale", [
                        "The retain threshold was raised from 2/3 ≈ 0.667 to 0.80 because "
                        "the xlsx classifies Overall Risk Score = 75 as Investigate (not "
                        "Retain) and ≥ 90 as Retain. Without the raise, Score-75 tx would "
                        "mis-route to Retain and disappear from the C&T queue.",
                    ], fill=SHORTCUT_BG, edge=SHORTCUT_EDGE, title_color=SHORTCUT_EDGE)


def slide_aggregation(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_header(s, "Contribution shares & case-level averaging",
                  "How the per-engine and overall figures in the C&T UI are computed")

    # Left: contribution share
    _add_bullet_box(s, Inches(0.5), Inches(1.3), Inches(6.2), Inches(2.8),
                    "Per-engine CONTRIBUTION SHARE (UI)", [
                        "For each engine: weighted = engine_avg × ENGINE_WEIGHTS[engine].",
                        "Total = sum of weighted contributions across the 4 engines.",
                        "Share displayed in 'Risk Signals' panel = weighted / total × 100 %.",
                        "By construction the four shares sum to ≈ 100 %, so officers can "
                        "eyeball which engine dominates the case.",
                        "Raw engine severity is also shown underneath the share (the band "
                        "that set the pre-baked value — 0.02–0.08 quiet, 0.55–0.70 firing).",
                    ], fill=WHITE, edge=EU_BLUE)

    # Right: case-level running average
    _add_bullet_box(s, Inches(7.0), Inches(1.3), Inches(5.8), Inches(2.8),
                    "Case-level RUNNING AVERAGE (backend)", [
                        "When a new order joins an existing case (C&T factory detects a "
                        "similar open case by Jaccard ≥ 0.4 on description):",
                        "  n = orders in the case after this one is appended",
                        "  Engine_X_new  = (Engine_X_old × (n-1) + order_Engine_X) / n",
                        "  Overall_new   = (Overall_old  × (n-1) + order_score)   / n",
                        "Each engine field on Sales_Order_Case is an arithmetic mean "
                        "across every order in the case — not a snapshot of the first "
                        "transaction.",
                    ], fill=WHITE, edge=EU_BLUE)

    # Caveat block at the bottom
    _add_bullet_box(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.8),
                    "⚠ Non-linearity caveat — shares are an approximation", [
                        "Per-order score passes through the weighted sum AND the vat_ratio "
                        "floor AND the min(1.0, …) cap. Only the linear part (weighted sum) "
                        "commutes with averaging.",
                        "Consequence: Overall_Case_Risk_Score ≠ Σ ENGINE_WEIGHTS × Engine_X_avg "
                        "whenever the floor or the cap fired on any order in the case.",
                        "For cases where no order ever crossed the floor (raw vat_ratio < 0.30 "
                        "everywhere) and no order ever hit the cap (total weighted ≤ 1.0), "
                        "the two are exactly equal and the shares reconstruct the overall.",
                        "For cases where floor/cap fired on some orders, the displayed "
                        "shares still sum to 100 % and still rank engines correctly by "
                        "contribution, but reconstructing the overall from them involves "
                        "a small residual that carries the floor/cap effect.",
                    ], fill=SHORTCUT_BG, edge=SHORTCUT_EDGE, title_color=SHORTCUT_EDGE,
                    body_size=12)


# ── Main ───────────────────────────────────────────────────────────────────

def main() -> None:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_overview(prs)
    slide_engine_vat_ratio(prs)
    slide_engine_ml(prs)
    slide_engine_ie(prs)
    slide_engine_vagueness(prs)
    slide_consolidation(prs)
    slide_thresholds(prs)
    slide_aggregation(prs)

    out = Path(__file__).resolve().parent.parent / "docs" / "risk_engines.pptx"
    prs.save(out)
    print(f"Wrote {out} ({out.stat().st_size:,} bytes, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
