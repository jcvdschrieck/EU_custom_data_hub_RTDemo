#!/usr/bin/env python3
"""
Generate EU_VAT_Audit_Swimlane.pptx

Two horizontal swim lanes:
  - EU VAT Hub (top)
  - Ireland App  (bottom)

Seven process steps flowing left-to-right, with annotated cross-lane arrows.
"""
from __future__ import annotations
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn as _qn
from lxml import etree

# ── Palette ───────────────────────────────────────────────────────────────────
EU_BLUE     = RGBColor(0x00, 0x33, 0x99)
EU_BLUE_LT  = RGBColor(0xCC, 0xD9, 0xF0)
EU_YELLOW   = RGBColor(0xFF, 0xED, 0x00)
IE_NAVY     = RGBColor(0x14, 0x5B, 0x94)
IE_NAVY_LT  = RGBColor(0xD0, 0xE8, 0xF5)
IE_TEAL     = RGBColor(0x00, 0x9E, 0xAB)
LLM_PURPLE  = RGBColor(0x6A, 0x3D, 0x9A)
DB_GREEN    = RGBColor(0x21, 0x7A, 0x3C)
DB_GREEN_LT = RGBColor(0xD0, 0xED, 0xD9)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xF4, 0xF4, 0xF4)
MID_GREY    = RGBColor(0xBB, 0xBB, 0xBB)
DARK_GREY   = RGBColor(0x33, 0x33, 0x33)


# ── Primitive helpers (all coordinates in inches, like the arch script) ───────

def _rect(slide, x, y, w, h,
          fill=WHITE, border=MID_GREY, border_pt=1.0,
          text="", font_size=10, bold=False,
          text_color=DARK_GREY, align=PP_ALIGN.CENTER,
          v_anchor=MSO_ANCHOR.MIDDLE, wrap=True):
    shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(border_pt)
    if text:
        tf = shp.text_frame
        tf.word_wrap = wrap
        tf.vertical_anchor = v_anchor
        p = tf.paragraphs[0]
        p.alignment = align
        for chunk in text.split("\n"):
            if p.runs:
                p = tf.add_paragraph()
                p.alignment = align
            run = p.add_run()
            run.text = chunk
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = text_color
    return shp


def _txt(slide, text, x, y, w, h,
         font_size=10, bold=False, color=DARK_GREY,
         align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    first = True
    for chunk in text.split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = chunk
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def _line(slide, x1, y1, x2, y2,
          color=DARK_GREY, width_pt=1.5,
          arrow_end=True, dashed=False):
    """Straight connector; arrow points at the destination end (headEnd)."""
    cx = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    cx.line.color.rgb = color
    cx.line.width = Pt(width_pt)
    ln = cx.line._ln
    if dashed:
        prstDash = etree.SubElement(ln, _qn("a:prstDash"))
        prstDash.set("val", "dash")
    if arrow_end:
        tail = ln.find(_qn("a:tailEnd"))
        if tail is None:
            tail = etree.SubElement(ln, _qn("a:tailEnd"))
        tail.set("type", "triangle")
        tail.set("w", "med")
        tail.set("len", "med")
    return cx


# ── Compound helpers ──────────────────────────────────────────────────────────

def _step_card(slide, x, y, w, h,
               step_num, title, body,
               hdr_fill, hdr_text=WHITE):
    """Numbered step card: coloured header + white body with wrapped body text."""
    HDR = 0.46   # header height in inches
    PAD = 0.11   # horizontal text padding

    # Card outline (full card, white body)
    _rect(slide, x, y, w, h, fill=WHITE, border=hdr_fill, border_pt=1.5)

    # Header band
    _rect(slide, x, y, w, HDR, fill=hdr_fill, border=hdr_fill, border_pt=0)

    # Badge — step number
    BADGE = 0.30
    bx = x + 0.10
    by = y + (HDR - BADGE) / 2
    _rect(slide, bx, by, BADGE, BADGE,
          fill=WHITE, border=hdr_fill, border_pt=0,
          text=str(step_num), font_size=10, bold=True,
          text_color=hdr_fill, align=PP_ALIGN.CENTER,
          v_anchor=MSO_ANCHOR.MIDDLE)

    # Header title
    _txt(slide, title,
         x + BADGE + 0.24, y + 0.06,
         w - BADGE - 0.34, HDR - 0.08,
         font_size=9, bold=True, color=hdr_text)

    # Body text
    _txt(slide, body,
         x + PAD, y + HDR + 0.09,
         w - PAD * 2, h - HDR - 0.12,
         font_size=8, color=DARK_GREY, align=PP_ALIGN.LEFT)


def _db_box(slide, x, y, w, h, label):
    """Simple database representation: dark green box with label."""
    _rect(slide, x, y, w, h,
          fill=DB_GREEN, border=DB_GREEN, border_pt=0,
          text=label, font_size=8, bold=True,
          text_color=WHITE, align=PP_ALIGN.CENTER,
          v_anchor=MSO_ANCHOR.MIDDLE)
    # Small cap to suggest cylinder top
    _rect(slide, x, y, w, h * 0.22,
          fill=DB_GREEN_LT, border=DB_GREEN, border_pt=0)


# ── Build presentation ────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

    # Background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_GREY

    # ── Title bar ─────────────────────────────────────────────────────────────
    TITLE_H = 0.55
    _rect(slide, 0, 0, 13.33, TITLE_H,
          fill=EU_BLUE, border=EU_BLUE,
          text="EU VAT Audit System — End-to-End Process Flow",
          font_size=16, bold=True, text_color=WHITE,
          align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.MIDDLE)

    # ── Lane dimensions ───────────────────────────────────────────────────────
    LANE_TOP   = TITLE_H + 0.08
    LABEL_W    = 1.22
    LANE_H     = (7.5 - LANE_TOP - 0.12) / 2   # ≈ 3.38
    EU_LANE_Y  = LANE_TOP
    IE_LANE_Y  = LANE_TOP + LANE_H

    # EU lane background
    _rect(slide, 0, EU_LANE_Y, 13.33, LANE_H,
          fill=EU_BLUE_LT, border=EU_BLUE_LT)
    # IE lane background
    _rect(slide, 0, IE_LANE_Y, 13.33, LANE_H,
          fill=IE_NAVY_LT, border=IE_NAVY_LT)

    # Lane separator
    _rect(slide, 0, IE_LANE_Y, 13.33, 0.02,
          fill=MID_GREY, border=MID_GREY)

    # EU label
    _rect(slide, 0, EU_LANE_Y, LABEL_W, LANE_H,
          fill=EU_BLUE, border=EU_BLUE,
          text="EU VAT Hub\n\nEuropean Institution\n(port 8503 · 8502)",
          font_size=9, bold=False, text_color=WHITE,
          align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)

    # IE label
    _rect(slide, 0, IE_LANE_Y, LABEL_W, LANE_H,
          fill=IE_NAVY, border=IE_NAVY,
          text="Ireland App\n\nIrish Revenue\n(port 8501)",
          font_size=9, bold=False, text_color=WHITE,
          align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)

    # ── Step card geometry ────────────────────────────────────────────────────
    CONTENT_X = LABEL_W + 0.12
    CONTENT_W = 13.33 - CONTENT_X - 0.08
    N         = 7
    GAP       = 0.11
    CARD_W    = (CONTENT_W - (N - 1) * GAP) / N   # ≈ 1.60
    PAD_V     = 0.18
    EU_CARD_Y = EU_LANE_Y + PAD_V
    IE_CARD_Y = IE_LANE_Y + PAD_V
    CARD_H    = LANE_H - 2 * PAD_V                # ≈ 3.02

    def cx(i):   # left edge of step i (0-based)
        return CONTENT_X + i * (CARD_W + GAP)

    # ── STEP 1 — EU: Centralise Invoice Records ───────────────────────────────
    _step_card(slide, cx(0), EU_CARD_Y, CARD_W, CARD_H,
               step_num=1,
               title="Centralise Invoice Records",
               body=("Synthetic dataset ingested: ~2,800 invoices across "
                     "10 EU member states.\n\n"
                     "Stores factual data only: parties, amounts, applied "
                     "VAT rates, transaction classification.\n\n"
                     "No risk scoring — that is each member state's "
                     "responsibility."),
               hdr_fill=EU_BLUE)

    # DB icon inside step 1 card (bottom area)
    DB_W = CARD_W * 0.62
    DB_H = 0.52
    DB_X = cx(0) + (CARD_W - DB_W) / 2
    DB_Y = EU_CARD_Y + CARD_H - DB_H - 0.10
    _db_box(slide, DB_X, DB_Y, DB_W, DB_H, "eu_vat.db")

    # ── STEP 2 — EU: REST API ─────────────────────────────────────────────────
    _step_card(slide, cx(1), EU_CARD_Y, CARD_W, CARD_H,
               step_num=2,
               title="Serve via REST API",
               body=("FastAPI (port 8503) exposes invoice data.\n\n"
                     "Key endpoint:\nGET /api/v1/invoices\n"
                     "  ?country=IE\n  &date_from=…\n\n"
                     "All inbound requests logged: timestamp, client "
                     "country, latency, records returned."),
               hdr_fill=EU_BLUE)

    # ── STEP 3 — IE: Fetch Increment ─────────────────────────────────────────
    _step_card(slide, cx(2), IE_CARD_Y, CARD_W, CARD_H,
               step_num=3,
               title="Fetch Increment",
               body=("Irish app queries EU Hub for invoices dated after "
                     "the Irish DB cutoff (> 2026-03-25).\n\n"
                     "HTTP GET with header:\nX-Client-Country: IE\n\n"
                     "Returns 25 new Irish records "
                     "(2026-03-26 → 2026-03-30).\n\n"
                     "Outbound calls logged to eu_query_log.db."),
               hdr_fill=IE_NAVY)

    # ── STEP 4 — IE: Pre-Classify by Risk ────────────────────────────────────
    _step_card(slide, cx(3), IE_CARD_Y, CARD_W, CARD_H,
               step_num=4,
               title="Pre-Classify by Risk",
               body=("Each invoice ranked using local supplier history "
                     "(vat_audit.db):\n\n"
                     "HIGH\n"
                     "  error rate ≥ 50%\n"
                     "  or gross > €15 000\n\n"
                     "MEDIUM\n"
                     "  error rate ≥ 15%\n"
                     "  or new supplier\n\n"
                     "LOW\n"
                     "  clean history"),
               hdr_fill=IE_NAVY)

    # ── STEP 5 — IE: Auditor Selects & Queues ────────────────────────────────
    _step_card(slide, cx(4), IE_CARD_Y, CARD_W, CARD_H,
               step_num=5,
               title="Auditor Selects & Queues",
               body=("Invoices displayed sorted:\nHIGH → MEDIUM → LOW.\n\n"
                     "Auditor ticks the invoices of interest.\n\n"
                     "On 'Launch VAT Analysis', a second EU Hub query "
                     "is made per selected record to retrieve the full "
                     "InvoiceDetail (including line items).\n\n"
                     "Step 3 returned summaries only — line items are "
                     "required for per-line LLM analysis in step 6."),
               hdr_fill=IE_NAVY)

    # ── STEP 6 — IE: LLM Compliance Analysis ─────────────────────────────────
    _step_card(slide, cx(5), IE_CARD_Y, CARD_W, CARD_H,
               step_num=6,
               title="LLM Compliance Analysis",
               body=("Per invoice:\n"
                     "1. RAG retrieves Irish VAT legislation from ChromaDB\n\n"
                     "2. LM Studio (port 1234) returns verdict per line "
                     "item:\n   correct / incorrect / uncertain\n\n"
                     "3. Overall verdict:\n"
                     "   any incorrect → incorrect\n\n"
                     "Every call logged to analysis_log.db."),
               hdr_fill=IE_NAVY)

    # LM Studio badge
    LB_W = CARD_W - 0.18
    LB_H = 0.28
    LB_X = cx(5) + 0.09
    LB_Y = IE_CARD_Y + CARD_H - LB_H - 0.08
    _rect(slide, LB_X, LB_Y, LB_W, LB_H,
          fill=LLM_PURPLE, border=LLM_PURPLE,
          text="LM Studio  (port 1234)",
          font_size=7, bold=True, text_color=WHITE,
          align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)

    # ── STEP 7 — IE: Save & Dashboard ────────────────────────────────────────
    _step_card(slide, cx(6), IE_CARD_Y, CARD_W, CARD_H,
               step_num=7,
               title="Save & Review in Dashboard",
               body=("Results persisted to:\n"
                     "  vat_audit.db\n"
                     "  history.json\n\n"
                     "Prioritization Dashboard ranks all analysed invoices "
                     "by risk score.\n\n"
                     "Auditor reviews verdicts, rationales, and legislation "
                     "excerpts.\n\n"
                     "Re-analysis available from History page."),
               hdr_fill=IE_NAVY)

    # ── Within-lane process arrows ────────────────────────────────────────────
    EU_MID = EU_LANE_Y + LANE_H / 2
    IE_MID = IE_LANE_Y + LANE_H / 2

    # Step 1 → 2 (EU lane)
    _line(slide, cx(0) + CARD_W, EU_MID, cx(1), EU_MID,
          color=EU_BLUE, width_pt=2.0)

    # Steps 3→4, 4→5, 5→6, 6→7 (IE lane)
    for i in [2, 3, 4, 5]:
        _line(slide, cx(i) + CARD_W, IE_MID, cx(i + 1), IE_MID,
              color=IE_NAVY, width_pt=2.0)

    # ── Cross-lane arrows: increment fetch (between steps 2 and 3) ────────────
    LANE_SEP = IE_LANE_Y   # y of the separator line

    EU_CARD_BTM = EU_CARD_Y + CARD_H   # bottom of EU cards
    IE_CARD_TOP = IE_CARD_Y            # top of IE cards

    # Choose x positions just to the right of step 2 / left of step 3
    OUT_X = cx(2) + CARD_W * 0.28   # outbound (IE → EU): dashed
    IN_X  = cx(2) + CARD_W * 0.68   # inbound  (EU → IE): solid

    # Outbound: IE sends request upward to EU
    _line(slide, OUT_X, IE_CARD_TOP, OUT_X, EU_CARD_BTM,
          color=IE_TEAL, width_pt=1.5, dashed=True)

    # Inbound: EU returns data downward to IE
    _line(slide, IN_X, EU_CARD_BTM, IN_X, IE_CARD_TOP,
          color=EU_BLUE, width_pt=1.5)

    # Labels in the gap between lanes
    GAP_Y = LANE_SEP - 0.50
    _txt(slide,
         "GET /api/v1/invoices\n?date_from=2026-03-26\nX-Client-Country: IE",
         OUT_X - 0.85, GAP_Y, 1.2, 0.50,
         font_size=6.5, color=IE_TEAL, align=PP_ALIGN.RIGHT)
    _txt(slide,
         "25 IE invoices (JSON)",
         IN_X + 0.04, GAP_Y + 0.10, 1.1, 0.30,
         font_size=6.5, color=EU_BLUE, align=PP_ALIGN.LEFT)

    # ── Cross-lane arrows: detail fetch (step 5, per selected invoice) ────────
    OUT5_X = cx(4) + CARD_W * 0.30
    IN5_X  = cx(4) + CARD_W * 0.70

    _line(slide, OUT5_X, IE_CARD_TOP, OUT5_X, EU_CARD_BTM,
          color=IE_TEAL, width_pt=1.5, dashed=True)
    _line(slide, IN5_X, EU_CARD_BTM, IN5_X, IE_CARD_TOP,
          color=EU_BLUE, width_pt=1.5)

    GAP5_Y = LANE_SEP - 0.44
    _txt(slide,
         "GET /api/v1/invoices/{id}",
         OUT5_X - 0.90, GAP5_Y, 1.1, 0.30,
         font_size=6.5, color=IE_TEAL, align=PP_ALIGN.RIGHT)
    _txt(slide,
         "InvoiceDetail\n(with line items)",
         IN5_X + 0.04, GAP5_Y, 1.0, 0.38,
         font_size=6.5, color=EU_BLUE, align=PP_ALIGN.LEFT)

    # ── Legend ────────────────────────────────────────────────────────────────
    LEG_Y = 7.5 - 0.26
    _txt(slide,
         "─── request / response within lane     "
         "- - - HTTP request (Ireland → EU Hub)     "
         "─── HTTP response (EU Hub → Ireland)",
         0.2, LEG_Y, 13.0, 0.24,
         font_size=7, color=MID_GREY, align=PP_ALIGN.LEFT)

    # ── Save ──────────────────────────────────────────────────────────────────
    out = "EU_VAT_Audit_Swimlane.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    build()
