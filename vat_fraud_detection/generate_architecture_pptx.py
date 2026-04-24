#!/usr/bin/env python3
"""Generate EU VAT Audit System — Architecture slides (PPTX)."""
from __future__ import annotations
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x1B, 0x2A, 0x4A)
IE_DARK     = RGBColor(0x1F, 0x49, 0x9C)
IE_MED      = RGBColor(0x44, 0x72, 0xC4)
IE_LIGHT    = RGBColor(0xD9, 0xE8, 0xF5)
EU_DARK     = RGBColor(0x1E, 0x6B, 0x4A)
EU_MED      = RGBColor(0x2F, 0x9E, 0x6F)
EU_LIGHT    = RGBColor(0xD5, 0xEF, 0xE3)
LM_DARK     = RGBColor(0x5E, 0x2A, 0x9A)
LM_MED      = RGBColor(0x7E, 0x57, 0xC2)
LM_LIGHT    = RGBColor(0xEB, 0xD9, 0xF7)
DB_DARK     = RGBColor(0xC4, 0x5E, 0x00)
DB_MED      = RGBColor(0xF4, 0xA0, 0x21)
DB_LIGHT    = RGBColor(0xFE, 0xF3, 0xDC)
GREY_DARK   = RGBColor(0x33, 0x33, 0x33)
GREY_MED    = RGBColor(0x88, 0x88, 0x88)
GREY_LIGHT  = RGBColor(0xF2, 0xF2, 0xF2)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
RED_DARK    = RGBColor(0xC0, 0x20, 0x20)


# ── Helpers ───────────────────────────────────────────────────────────────────

def _set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, x, y, w, h,
         fill=WHITE, border=GREY_MED, border_pt=1.0,
         text="", font_size=10, bold=False,
         text_color=BLACK, align=PP_ALIGN.CENTER,
         v_anchor=MSO_ANCHOR.MIDDLE, wrap=True):
    shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(border_pt)
    if text:
        tf = shp.text_frame
        tf.word_wrap = wrap
        tf.vertical_anchor = v_anchor
        # clear default paragraph
        p = tf.paragraphs[0]
        p.alignment = align
        for line in text.split("\n"):
            if p.runs:
                p = tf.add_paragraph()
                p.alignment = align
            run = p.add_run()
            run.text = line
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = text_color
    return shp


def txtbox(slide, text, x, y, w, h,
           font_size=10, bold=False, color=BLACK,
           align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def multiline_txtbox(slide, lines: list[tuple], x, y, w, h,
                     v_anchor=MSO_ANCHOR.TOP):
    """lines = list of (text, font_size, bold, color, align)"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = v_anchor
    first = True
    for (text, fs, bold, color, align) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(fs)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def line(slide, x1, y1, x2, y2,
         color=GREY_DARK, width_pt=1.5,
         arrow_end=True, arrow_start=False, dashed=False):
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE
    from pptx.oxml.ns import qn as _qn
    cx = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    cx.line.color.rgb = color
    cx.line.width = Pt(width_pt)
    ln = cx.line._ln
    if dashed:
        prstDash = etree.SubElement(ln, _qn('a:prstDash'))
        prstDash.set('val', 'dash')
    if arrow_end:
        head = ln.find(_qn('a:headEnd'))
        if head is None:
            head = etree.SubElement(ln, _qn('a:headEnd'))
        head.set('type', 'arrow')
        head.set('w', 'med')
        head.set('len', 'med')
    if arrow_start:
        tail = ln.find(_qn('a:tailEnd'))
        if tail is None:
            tail = etree.SubElement(ln, _qn('a:tailEnd'))
        tail.set('type', 'arrow')
        tail.set('w', 'med')
        tail.set('len', 'med')
    return cx


def slide_title(slide, title: str, subtitle: str = ""):
    rect(slide, 0, 0, 13.33, 0.65,
         fill=NAVY, border=NAVY,
         text=title, font_size=18, bold=True,
         text_color=WHITE, v_anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        txtbox(slide, subtitle, 0.15, 0.67, 13.0, 0.35,
               font_size=9, color=GREY_MED, italic=True)


def new_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, WHITE)
    return slide


# ── Presentation ──────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════
s1 = new_slide(prs)
_set_bg(s1, NAVY)

rect(s1, 1.5, 1.8, 10.33, 1.5,
     fill=NAVY, border=IE_MED, border_pt=2,
     text="EU VAT Audit System", font_size=40, bold=True,
     text_color=WHITE)
txtbox(s1, "Architecture Overview", 1.5, 3.4, 10.33, 0.6,
       font_size=20, color=IE_LIGHT, align=PP_ALIGN.CENTER)
txtbox(s1, "Functional  ·  Technical  ·  Data & Request Flows",
       1.5, 4.1, 10.33, 0.5,
       font_size=13, color=GREY_MED, align=PP_ALIGN.CENTER, italic=True)

for label_txt, col, xc in [
    ("🇮🇪 Ireland VAT App\nport 8501", IE_MED, 2.2),
    ("🏛️ EU VAT Hub\nAPI 8503 · UI 8502", EU_MED, 6.16),
    ("⚙️ LM Studio\nport 1234", LM_MED, 10.13),
]:
    rect(s1, xc, 5.2, 2.8, 1.0,
         fill=RGBColor(0x25, 0x35, 0x55), border=col, border_pt=1.5,
         text=label_txt, font_size=11, text_color=WHITE)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Functional Architecture
# ═══════════════════════════════════════════════════════════════
s2 = new_slide(prs)
slide_title(s2, "Functional Architecture",
            "Actors, systems and their responsibilities")

# ── Ireland VAT App box ──────────────────────────────────────
rect(s2, 0.25, 1.1, 4.3, 5.5,
     fill=IE_LIGHT, border=IE_DARK, border_pt=2)
rect(s2, 0.35, 1.1, 4.1, 0.45,
     fill=IE_DARK, border=IE_DARK,
     text="🇮🇪  Ireland VAT App  (port 8501)",
     font_size=11, bold=True, text_color=WHITE)

ie_bullets = [
    "· Maintains local Irish invoice DB (cutoff: 2026-03-25)",
    "· Queries EU Hub for post-cutoff increment",
    "· Pre-classifies invoices HIGH / MEDIUM / LOW",
    "· Queues invoices for LLM compliance analysis",
    "· Verdict per line item: correct / incorrect / uncertain",
    "· Prioritization Dashboard with risk scoring",
    "· Activity log for LLM calls & outbound requests",
]
for i, b in enumerate(ie_bullets):
    txtbox(s2, b, 0.45, 1.65 + i * 0.44, 4.0, 0.42,
           font_size=9.5, color=IE_DARK)

# ── EU VAT Hub box ───────────────────────────────────────────
rect(s2, 5.35, 1.1, 4.3, 5.5,
     fill=EU_LIGHT, border=EU_DARK, border_pt=2)
rect(s2, 5.45, 1.1, 4.1, 0.45,
     fill=EU_DARK, border=EU_DARK,
     text="🏛️  EU VAT Hub  (API 8503 · UI 8502)",
     font_size=11, bold=True, text_color=WHITE)

eu_bullets = [
    "· Central invoice repository for all member states",
    "· Stores factual data only: amounts, parties,",
    "   VAT rates applied, transaction classification",
    "· No risk scoring — assessment is country's role",
    "· ~2,800 synthetic invoices across 10 countries",
    "· Logs all inbound API requests (with country)",
    "· Read-only EU Administrator dashboard",
]
for i, b in enumerate(eu_bullets):
    txtbox(s2, b, 5.45, 1.65 + i * 0.44, 4.15, 0.42,
           font_size=9.5, color=EU_DARK)

# ── LM Studio box ────────────────────────────────────────────
rect(s2, 1.0, 7.0 - 1.7, 2.6, 1.55,  # ~y=5.3
     fill=LM_LIGHT, border=LM_DARK, border_pt=1.5)
rect(s2, 1.0, 5.3, 2.6, 0.38,
     fill=LM_DARK, border=LM_DARK,
     text="⚙️  LM Studio  (port 1234)",
     font_size=10, bold=True, text_color=WHITE)
txtbox(s2,
       "· Chat model — VAT analysis\n· Embedding model — legislation RAG\n· OpenAI-compatible local API",
       1.05, 5.72, 2.5, 1.0, font_size=9, color=LM_DARK)

# ── ChromaDB box ─────────────────────────────────────────────
rect(s2, 1.0 + 2.85, 5.3, 2.2, 1.55,
     fill=DB_LIGHT, border=DB_DARK, border_pt=1.5)
rect(s2, 3.85, 5.3, 2.2, 0.38,
     fill=DB_DARK, border=DB_DARK,
     text="🗄️  ChromaDB",
     font_size=10, bold=True, text_color=WHITE)
txtbox(s2,
       "· VAT legislation chunks\n· Retrieved per line item\n· Local vector store",
       3.9, 5.72, 2.1, 1.0, font_size=9, color=DB_DARK)

# ── Irish Auditor ────────────────────────────────────────────
rect(s2, 9.9, 2.6, 1.6, 0.9,
     fill=GREY_LIGHT, border=GREY_DARK, border_pt=1.5,
     text="👤\nIrish Auditor", font_size=10, bold=True, text_color=GREY_DARK)
# EU Admin
rect(s2, 11.7, 2.6, 1.5, 0.9,
     fill=GREY_LIGHT, border=GREY_DARK, border_pt=1.5,
     text="👤\nEU Admin", font_size=10, bold=True, text_color=GREY_DARK)

# ── Arrows ───────────────────────────────────────────────────
# Ireland ↔ EU Hub (HTTP REST)
line(s2, 4.55, 3.85, 5.35, 3.85,
     color=GREY_DARK, width_pt=2.0, arrow_end=True, arrow_start=True)
txtbox(s2, "HTTP REST\nX-Client-Country: IE",
       4.5, 3.55, 1.0, 0.6, font_size=7.5, color=GREY_DARK, align=PP_ALIGN.CENTER)

# Auditor → Ireland
line(s2, 9.9, 3.05, 9.65, 3.05,
     color=IE_MED, width_pt=1.5, arrow_end=False, arrow_start=True)
# Auditor → EU Hub (to the box right edge area)
line(s2, 9.9, 3.55, 9.65, 3.55,
     color=EU_MED, width_pt=1.5, arrow_end=False, arrow_start=True)
# EU Admin → EU Hub
line(s2, 11.7, 3.05, 9.65, 3.05,
     color=EU_MED, width_pt=1.5, arrow_end=True, arrow_start=False)

# Ireland → LM Studio
line(s2, 2.5, 6.6, 2.5, 6.85,
     color=LM_MED, width_pt=1.5, arrow_end=True)
txtbox(s2, "LLM calls\n(local)", 2.55, 6.6, 1.0, 0.35,
       font_size=7.5, color=LM_MED)

# Ireland → ChromaDB
line(s2, 3.2, 6.6, 4.3, 6.0,
     color=DB_MED, width_pt=1.5, arrow_end=True)
txtbox(s2, "RAG", 3.5, 6.22, 0.6, 0.3,
       font_size=7.5, color=DB_MED)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Technical Architecture Overview
# ═══════════════════════════════════════════════════════════════
s3 = new_slide(prs)
slide_title(s3, "Technical Architecture — Overview",
            "All components, entry points and connections at a glance")

# ── Ireland App group ────────────────────────────────────────
rect(s3, 0.15, 1.05, 5.8, 5.85,
     fill=IE_LIGHT, border=IE_DARK, border_pt=1.5)
rect(s3, 0.15, 1.05, 5.8, 0.38,
     fill=IE_DARK, border=IE_DARK,
     text="🇮🇪  Ireland VAT App  (Streamlit · port 8501)",
     font_size=10, bold=True, text_color=WHITE)

# Pages sub-box
rect(s3, 0.25, 1.55, 2.65, 3.0,
     fill=WHITE, border=IE_MED, border_pt=1.0,
     text="pages/", font_size=9, bold=True,
     text_color=IE_DARK, v_anchor=MSO_ANCHOR.TOP)
for i, pg in enumerate([
    "1 · Invoice Analyzer",
    "2 · Prioritization Dashboard",
    "3 · Case View",
    "4 · History",
    "5 · EU Query / Increment",
    "6 · Activity Log",
]):
    txtbox(s3, pg, 0.35, 1.92 + i * 0.37, 2.45, 0.35,
           font_size=8.5, color=GREY_DARK)

# Lib sub-box
rect(s3, 3.05, 1.55, 2.8, 3.0,
     fill=WHITE, border=IE_MED, border_pt=1.0,
     text="lib/", font_size=9, bold=True,
     text_color=IE_DARK, v_anchor=MSO_ANCHOR.TOP)
for i, mod in enumerate([
    "analyser.py",
    "analysis_log.py",
    "eu_client.py",
    "database.py  ·  db_seeder.py",
    "persistence.py",
    "rag.py  ·  vector_store.py",
]):
    txtbox(s3, mod, 3.15, 1.92 + i * 0.37, 2.6, 0.35,
           font_size=8.5, color=GREY_DARK)

# Data sub-box
rect(s3, 0.25, 4.65, 5.6, 2.1,
     fill=DB_LIGHT, border=DB_MED, border_pt=1.0,
     text="data/", font_size=9, bold=True,
     text_color=DB_DARK, v_anchor=MSO_ANCHOR.TOP)
ie_data = [
    ("vat_audit.db", 0.4),
    ("analysis_log.db", 2.1),
    ("eu_query_log.db", 3.8),
    ("history.json", 0.4),
    ("chroma_db/", 2.1),
]
for name, xoff in ie_data:
    rect(s3, 0.35 + xoff, 5.05, 1.55, 0.38,
         fill=DB_LIGHT, border=DB_DARK, border_pt=0.5,
         text=name, font_size=7.5, text_color=DB_DARK)

# ── EU Hub group ────────────────────────────────────────────
rect(s3, 6.3, 1.05, 6.85, 5.85,
     fill=EU_LIGHT, border=EU_DARK, border_pt=1.5)
rect(s3, 6.3, 1.05, 6.85, 0.38,
     fill=EU_DARK, border=EU_DARK,
     text="🏛️  EU VAT Hub",
     font_size=10, bold=True, text_color=WHITE)

# API sub-box
rect(s3, 6.45, 1.55, 3.15, 3.5,
     fill=WHITE, border=EU_MED, border_pt=1.0,
     text="api.py  (FastAPI · port 8503)", font_size=9, bold=True,
     text_color=EU_DARK, v_anchor=MSO_ANCHOR.TOP)
for i, ep in enumerate([
    "GET /health",
    "GET /api/v1/invoices",
    "GET /api/v1/invoices/{id}",
    "GET /api/v1/stats/by-country",
    "GET /api/v1/stats/by-tx-type",
    "GET /api/v1/stats/by-treatment",
    "GET /api/v1/logs",
]):
    txtbox(s3, ep, 6.6, 1.92 + i * 0.43, 2.9, 0.4,
           font_size=8.5, color=GREY_DARK)

# Middleware
rect(s3, 6.45, 5.12, 3.15, 0.65,
     fill=EU_LIGHT, border=EU_MED, border_pt=1.0,
     text="logging_middleware.py\nCaptures: timestamp · client_country · latency · records",
     font_size=8.0, text_color=EU_DARK)

# Dashboard sub-box
rect(s3, 9.75, 1.55, 3.25, 3.5,
     fill=WHITE, border=EU_MED, border_pt=1.0,
     text="app.py  (Streamlit · port 8502)", font_size=9, bold=True,
     text_color=EU_DARK, v_anchor=MSO_ANCHOR.TOP)
for i, pg in enumerate([
    "1 · Overview",
    "2 · Invoice Browser",
    "3 · Analytics",
    "4 · Activity Log",
]):
    txtbox(s3, pg, 9.85, 1.92 + i * 0.45, 3.05, 0.4,
           font_size=8.5, color=GREY_DARK)

# seeder + lib
rect(s3, 9.75, 5.12, 3.25, 0.65,
     fill=EU_LIGHT, border=EU_MED, border_pt=1.0,
     text="lib/  seeder.py · database.py · models.py",
     font_size=8.0, text_color=EU_DARK)

# EU Data
rect(s3, 6.45, 5.85, 6.55, 0.85,
     fill=DB_LIGHT, border=DB_MED, border_pt=1.0,
     text="data/eu_vat.db  —  ~2,800 records · 10 member states  ·  tables: invoices · line_items · api_log",
     font_size=8.5, text_color=DB_DARK)

# ── External ────────────────────────────────────────────────
rect(s3, 0.15, 7.0 - 0.5, 2.8, 0.42,
     fill=LM_LIGHT, border=LM_DARK, border_pt=1.0,
     text="⚙️ LM Studio (port 1234)  Chat + Embedding",
     font_size=9, bold=True, text_color=LM_DARK)

# ── Connections ───────────────────────────────────────────────
# Ireland pages → lib
line(s3, 2.9, 2.5, 3.05, 2.5, color=IE_MED, width_pt=1.0, arrow_end=True)
# Ireland → EU Hub (HTTP)
line(s3, 6.1, 3.8, 6.3, 3.8,
     color=GREY_DARK, width_pt=2.0, arrow_end=True, arrow_start=True)
txtbox(s3, "HTTP\nREST", 6.08, 3.5, 0.8, 0.55,
       font_size=7.5, color=GREY_DARK, align=PP_ALIGN.CENTER)
# Ireland lib → LM Studio
line(s3, 2.0, 6.85, 2.0, 7.0 - 0.5,
     color=LM_MED, width_pt=1.5, arrow_end=True)
# Ireland lib → chroma (within data)
line(s3, 3.7, 4.65, 3.95, 4.3,
     color=DB_MED, width_pt=1.0, arrow_end=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Ireland App Component Detail
# ═══════════════════════════════════════════════════════════════
s4 = new_slide(prs)
slide_title(s4, "Technical Architecture — Ireland VAT App",
            "Pages, library modules, data stores and their wiring")

# Pages column
rect(s4, 0.2, 1.05, 4.1, 0.38,
     fill=IE_DARK, border=IE_DARK,
     text="pages/", font_size=10, bold=True, text_color=WHITE)
pages = [
    ("1 · Invoice Analyzer", "Analysis queue runner — processes queued invoices,\nrenders verdict table + full rationale expander."),
    ("2 · Prioritization Dashboard", "Reads vat_audit.db, ranks by risk score,\nperiod filters, lazy-loads rationale."),
    ("3 · Case View", "Single invoice deep-dive with verdict and\nlegislation references."),
    ("4 · History", "Multi-select re-analysis. Loads past results from\npersistence.py, queues chosen invoices."),
    ("5 · EU Query / Increment", "Filters EU Hub invoices. Increment tab fetches\nIE records post-cutoff, pre-classifies, launches analysis."),
    ("6 · Activity Log", "Displays analysis_log.db — every LLM call\nwith timestamp, verdict, latency."),
]
for i, (title, desc) in enumerate(pages):
    rect(s4, 0.2, 1.5 + i * 0.92, 4.1, 0.85,
         fill=IE_LIGHT, border=IE_MED, border_pt=1.0)
    txtbox(s4, title, 0.3, 1.53 + i * 0.92, 3.9, 0.32,
           font_size=9, bold=True, color=IE_DARK)
    txtbox(s4, desc, 0.3, 1.82 + i * 0.92, 3.9, 0.5,
           font_size=8, color=GREY_DARK)

# Lib column
rect(s4, 4.6, 1.05, 4.5, 0.38,
     fill=IE_DARK, border=IE_DARK,
     text="lib/", font_size=10, bold=True, text_color=WHITE)
mods = [
    ("analyser.py", "RAG retrieval + LM Studio LLM call. Parses verdict JSON.\nWrites timing to analysis_log on every call."),
    ("analysis_log.py", "SQLite log: timestamp, invoice, supplier, model,\nline count, verdict, response_time_ms."),
    ("eu_client.py", "httpx client → EU Hub API. Sets X-Client-Country: IE.\nLogs all outbound requests to eu_query_log.db."),
    ("database.py / db_seeder.py", "Irish invoice SQLite DB (vat_audit.db).\nSeeded with synthetic data, cutoff 2026-03-25."),
    ("persistence.py", "save_result() → appends history.json + syncs to\nvat_audit.db. load_results() hydrates AnalysisResult list."),
    ("rag.py / vector_store.py", "Per-line-item ChromaDB similarity search.\nDeduplicates & caps at 12 chunks for the prompt."),
],
for i, (name, desc) in enumerate(mods[0]):
    rect(s4, 4.6, 1.5 + i * 0.92, 4.5, 0.85,
         fill=IE_LIGHT, border=IE_MED, border_pt=1.0)
    txtbox(s4, name, 4.7, 1.53 + i * 0.92, 4.3, 0.32,
           font_size=9, bold=True, color=IE_DARK)
    txtbox(s4, desc, 4.7, 1.82 + i * 0.92, 4.3, 0.5,
           font_size=8, color=GREY_DARK)

# Data column
rect(s4, 9.4, 1.05, 3.75, 0.38,
     fill=DB_DARK, border=DB_DARK,
     text="data/", font_size=10, bold=True, text_color=WHITE)
dbs = [
    ("vat_audit.db", "SQLite. Tables: invoices, analyses.\nPrimary Irish audit database."),
    ("history.json", "JSON array of serialised AnalysisResult objects.\nLoaded by persistence.py."),
    ("analysis_log.db", "One row per LLM call.\nPowers page 6 — Activity Log."),
    ("eu_query_log.db", "One row per outbound HTTP request to EU Hub.\nShown in EU Query → Activity Log tab."),
    ("chroma_db/", "Local ChromaDB vector store.\nIndexed from Irish VAT legislation PDFs."),
]
for i, (name, desc) in enumerate(dbs):
    rect(s4, 9.4, 1.5 + i * 1.12, 3.75, 1.05,
         fill=DB_LIGHT, border=DB_MED, border_pt=1.0)
    txtbox(s4, name, 9.5, 1.53 + i * 1.12, 3.55, 0.32,
           font_size=9, bold=True, color=DB_DARK)
    txtbox(s4, desc, 9.5, 1.85 + i * 1.12, 3.55, 0.55,
           font_size=8, color=GREY_DARK)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — EU VAT Hub Component Detail
# ═══════════════════════════════════════════════════════════════
s5 = new_slide(prs)
slide_title(s5, "Technical Architecture — EU VAT Hub",
            "FastAPI REST API, Streamlit dashboard, middleware and data model")

# API box
rect(s5, 0.2, 1.05, 5.8, 0.38,
     fill=EU_DARK, border=EU_DARK,
     text="api.py  —  FastAPI REST API  (port 8503)",
     font_size=10, bold=True, text_color=WHITE)
endpoints = [
    ("GET /health", "Returns {status, db_records}. Used by Ireland app on startup."),
    ("GET /api/v1/invoices", "Filtered list. Params: country, date_from/to, tx_type,\nscope, vat_treatment, description, limit, offset."),
    ("GET /api/v1/invoices/{id}", "Full invoice detail including line_items[]. Returns 404 if not found."),
    ("GET /api/v1/stats/by-country", "Aggregate totals (net, VAT, gross, count) grouped by country."),
    ("GET /api/v1/stats/by-transaction-type", "Count and totals by B2B/B2C and domestic/intra_EU/extra_EU."),
    ("GET /api/v1/stats/by-vat-treatment", "Count and totals by treatment (standard, reduced, zero…)."),
    ("GET /api/v1/logs", "Returns inbound api_log rows. Limit param up to 1000."),
]
for i, (ep, desc) in enumerate(endpoints):
    rect(s5, 0.2, 1.5 + i * 0.74, 5.8, 0.68,
         fill=EU_LIGHT, border=EU_MED, border_pt=0.8)
    txtbox(s5, ep, 0.3, 1.53 + i * 0.74, 5.6, 0.3,
           font_size=9, bold=True, color=EU_DARK)
    txtbox(s5, desc, 0.3, 1.82 + i * 0.74, 5.6, 0.38,
           font_size=8, color=GREY_DARK)

# Middleware box
rect(s5, 0.2, 6.7, 5.8, 0.7,
     fill=RGBColor(0xC8, 0xE8, 0xD8), border=EU_MED, border_pt=1.0)
txtbox(s5, "logging_middleware.py  —  ApiLoggingMiddleware", 0.3, 6.72, 5.6, 0.3,
       font_size=9, bold=True, color=EU_DARK)
txtbox(s5, "Wraps every request. Records: timestamp (UTC), method, endpoint, client_country (from X-Client-Country), "
           "status_code, response_time_ms, records_returned → api_log table.",
       0.3, 6.98, 5.6, 0.38, font_size=8, color=GREY_DARK)

# Dashboard box
rect(s5, 6.35, 1.05, 3.0, 0.38,
     fill=EU_DARK, border=EU_DARK,
     text="app.py  —  Streamlit Dashboard  (port 8502)",
     font_size=10, bold=True, text_color=WHITE)
dash_pages = [
    ("1 · Overview", "KPIs: total records, countries, latest requests.\nRecent inbound activity feed."),
    ("2 · Invoice Browser", "Filterable invoice list with country, date,\ntransaction type, VAT treatment filters."),
    ("3 · Analytics", "Charts: volume by country, transaction type,\nVAT treatment distribution."),
    ("4 · Activity Log", "Inbound request log with formatted UTC timestamps,\nclient country flag, status, latency."),
]
for i, (pg, desc) in enumerate(dash_pages):
    rect(s5, 6.35, 1.5 + i * 1.42, 3.0, 1.3,
         fill=EU_LIGHT, border=EU_MED, border_pt=0.8)
    txtbox(s5, pg, 6.45, 1.53 + i * 1.42, 2.8, 0.35,
           font_size=9, bold=True, color=EU_DARK)
    txtbox(s5, desc, 6.45, 1.88 + i * 1.42, 2.8, 0.7,
           font_size=8.5, color=GREY_DARK)

# lib box
rect(s5, 9.55, 1.05, 3.6, 0.38,
     fill=EU_DARK, border=EU_DARK,
     text="lib/", font_size=10, bold=True, text_color=WHITE)
lib_items = [
    ("database.py", "init_db(), query_invoices(), count_invoices(),\nget_invoice(), get_line_items(), write_api_log(),\nget_api_logs(), stats_by_*()."),
    ("seeder.py", "seed_if_empty() called at startup. Generates ~2,800\nsynthetic invoices (10 countries, realistic VAT errors)\n+ 25 IE increment records dated 2026-03-26 to 2026-03-30."),
    ("models.py", "Pydantic models: InvoiceSummary, InvoiceDetail,\nLineItemSummary, ApiLogEntry, stats models.\nNo risk fields — factual data only."),
    ("logging_middleware.py", "Starlette BaseHTTPMiddleware. Measures elapsed time,\nextracts X-Client-Country, reads X-Records-Returned\nfrom response, writes to api_log."),
]
for i, (name, desc) in enumerate(lib_items):
    rect(s5, 9.55, 1.5 + i * 1.47, 3.6, 1.35,
         fill=EU_LIGHT, border=EU_MED, border_pt=0.8)
    txtbox(s5, name, 9.65, 1.53 + i * 1.47, 3.4, 0.32,
           font_size=9, bold=True, color=EU_DARK)
    txtbox(s5, desc, 9.65, 1.85 + i * 1.47, 3.4, 0.8,
           font_size=8, color=GREY_DARK)

# Data
rect(s5, 6.35, 6.88 - 0.85, 6.8, 0.9,
     fill=DB_LIGHT, border=DB_MED, border_pt=1.0)
txtbox(s5, "data/eu_vat.db", 6.45, 6.05, 6.6, 0.3,
       font_size=9, bold=True, color=DB_DARK)
txtbox(s5, "Tables: invoices (invoice_id, dates, parties, amounts, vat_rate_applied, tx_type, tx_scope, vat_treatment, reporting_country) · "
           "line_items (description, category, quantity, unit_price, vat_rate_applied, net/vat_amount) · "
           "api_log (inbound request audit trail)",
       6.45, 6.35, 6.6, 0.55, font_size=7.5, color=GREY_DARK)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — Data Flow: Increment Fetch & Analysis Queue
# ═══════════════════════════════════════════════════════════════
s6 = new_slide(prs)
slide_title(s6, "Data Flow 1 — Increment Fetch & Analysis Queue",
            "Irish auditor fetches post-cutoff invoices from the EU Hub and queues them for analysis")

# Participants
PARTICIPANTS = [
    ("👤 Auditor",        IE_DARK,  WHITE),
    ("Ireland\nApp",      IE_MED,   WHITE),
    ("Irish DB\nvat_audit.db", DB_MED, WHITE),
    ("eu_client\n.py",    IE_DARK,  WHITE),
    ("eu_query\n_log.db", DB_DARK,  WHITE),
    ("EU Hub\nAPI :8503", EU_DARK,  WHITE),
    ("api_log\neu_vat.db",EU_MED,   WHITE),
]

N = len(PARTICIPANTS)
LEFT_MARGIN = 0.15
USABLE_W = 13.0
COL_W = USABLE_W / N
PART_H = 0.72
PART_Y = 1.05
LINE_TOP = PART_Y + PART_H
LINE_BOT = 7.35
STEP_START_Y = LINE_TOP + 0.18
STEP_SPACING = 0.72

col_cx = [LEFT_MARGIN + (i + 0.5) * COL_W for i in range(N)]

for i, (name, fill, tc) in enumerate(PARTICIPANTS):
    rect(s6, LEFT_MARGIN + i * COL_W + 0.05, PART_Y,
         COL_W - 0.1, PART_H,
         fill=fill, border=fill,
         text=name, font_size=8.5, bold=True, text_color=tc)
    # lifeline
    line(s6, col_cx[i], LINE_TOP, col_cx[i], LINE_BOT,
         color=GREY_MED, width_pt=0.75, arrow_end=False, dashed=True)

def seq_arrow(slide, from_i, to_i, step_y, label, color=GREY_DARK, ret=False):
    x1, x2 = col_cx[from_i], col_cx[to_i]
    lc = color
    line(slide, x1, step_y, x2, step_y,
         color=lc, width_pt=1.5,
         arrow_end=(not ret), arrow_start=ret)
    mid_x = (x1 + x2) / 2 - 0.6
    txtbox(slide, label, mid_x, step_y - 0.26, 1.3, 0.28,
           font_size=7, color=lc, align=PP_ALIGN.CENTER)

def step_note(slide, step_y, text, x_note=0.17, w_note=None):
    # Small circled step number on side - just a label
    pass

STEPS_6 = [
    (0, 1, "Click 'Fetch Increment'",          GREY_DARK, False),
    (1, 3, "fetch_increment(limit=500)",        IE_MED,    False),
    (3, 5, "GET /api/v1/invoices\n?date_from=2026-03-26\nX-Client-Country: IE", EU_MED, False),
    (6, 5, "write_api_log(timestamp,IE,…)",     EU_MED,    True),
    (5, 3, "JSON {total, items[25]}\nX-Records-Returned: 25", EU_MED, True),
    (3, 4, "write_log(GET, 200, 25, latency)",  DB_MED,    False),
    (3, 1, "items[]",                           IE_MED,    True),
    (1, 2, "SELECT supplier stats",             DB_MED,    False),
    (2, 1, "supplier error rates",              DB_MED,    True),
    (1, 1, "pre-classify HIGH/MED/LOW",         IE_MED,    False),  # self
    (0, 1, "Tick invoices →\n'Launch VAT Analysis'", GREY_DARK, False),
    (1, 3, "get_invoice(id) ×N",               IE_MED,    False),
    (3, 5, "GET /api/v1/invoices/{id}\n×N",    EU_MED,    False),
    (5, 3, "InvoiceDetail (line_items[])",      EU_MED,    True),
    (1, 1, "→ analysis_queue\n→ Invoice Analyzer", IE_DARK, False),
]

for step_idx, (fi, ti, lbl, col, ret) in enumerate(STEPS_6):
    sy = STEP_START_Y + step_idx * STEP_SPACING
    if sy > LINE_BOT - 0.1:
        break
    if fi == ti:
        # self-reference: small box to the right
        cx = col_cx[fi]
        rect(s6, cx - 0.55, sy - 0.14, 1.1, 0.36,
             fill=IE_LIGHT, border=col, border_pt=0.8,
             text=lbl, font_size=6.5, text_color=col)
    else:
        seq_arrow(s6, fi, ti, sy, lbl, color=col, ret=ret)
    # step number
    txtbox(s6, str(step_idx + 1), 0.05, sy - 0.18, 0.2, 0.25,
           font_size=7, bold=True, color=GREY_MED)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Data Flow: LLM Analysis Pipeline
# ═══════════════════════════════════════════════════════════════
s7 = new_slide(prs)
slide_title(s7, "Data Flow 2 — LLM Compliance Analysis",
            "Per-invoice pipeline: RAG retrieval → LM Studio → verdict parsing → persistence")

PARTICIPANTS_7 = [
    ("👤 Auditor",       IE_DARK,  WHITE),
    ("Invoice\nAnalyzer", IE_MED,  WHITE),
    ("analyser\n.py",    IE_DARK,  WHITE),
    ("rag.py",           IE_MED,   WHITE),
    ("ChromaDB\nchroma_db/", DB_DARK, WHITE),
    ("LM Studio\n:1234", LM_DARK,  WHITE),
    ("analysis\n_log.db", DB_MED,  WHITE),
    ("persistence\n.py", IE_DARK,  WHITE),
    ("vat_audit\n.db",   DB_DARK,  WHITE),
]

N7 = len(PARTICIPANTS_7)
COL_W7 = USABLE_W / N7
col_cx7 = [LEFT_MARGIN + (i + 0.5) * COL_W7 for i in range(N7)]

for i, (name, fill, tc) in enumerate(PARTICIPANTS_7):
    rect(s7, LEFT_MARGIN + i * COL_W7 + 0.04, PART_Y,
         COL_W7 - 0.08, PART_H,
         fill=fill, border=fill,
         text=name, font_size=8, bold=True, text_color=tc)
    line(s7, col_cx7[i], LINE_TOP, col_cx7[i], LINE_BOT,
         color=GREY_MED, width_pt=0.75, arrow_end=False, dashed=True)

def seq_arrow7(slide, fi, ti, sy, lbl, col=GREY_DARK, ret=False):
    x1, x2 = col_cx7[fi], col_cx7[ti]
    line(slide, x1, sy, x2, sy,
         color=col, width_pt=1.5,
         arrow_end=(not ret), arrow_start=ret)
    mid_x = (min(x1, x2) + abs(x2 - x1) / 2) - 0.55
    txtbox(slide, lbl, mid_x, sy - 0.26, 1.2, 0.28,
           font_size=6.5, color=col, align=PP_ALIGN.CENTER)

STEPS_7 = [
    (0, 1, "▶ Run Analysis",                     GREY_DARK, False),
    (1, 2, "analyse(invoice)",                    IE_MED,    False),
    (2, 3, "retrieve(line_item)\n× each line",    IE_MED,    False),
    (3, 5, "POST /v1/embeddings\n(description + category)", LM_MED, False),
    (5, 3, "embedding vector",                    LM_MED,    True),
    (3, 4, "similarity_search(vector, top-k)",    DB_MED,    False),
    (4, 3, "legislation chunks[]",                DB_MED,    True),
    (2, 2, "dedup · cap 12 chunks\nformat_context()", IE_MED, False),
    (2, 5, "POST /v1/chat/completions\n(invoice JSON + legislation)", LM_MED, False),
    (5, 2, "JSON {verdicts[]}",                   LM_MED,    True),
    (2, 2, "_overall_verdict()\nany incorrect → incorrect\nall correct → correct", IE_DARK, False),
    (2, 6, "write_log(ts, invoice,\nverdict, latency_ms)", DB_MED, False),
    (2, 1, "AnalysisResult",                      IE_MED,    True),
    (1, 7, "save_result(result)",                 IE_MED,    False),
    (7, 8, "INSERT analyses + invoices",          DB_MED,    False),
    (1, 0, "render verdict table\n+ rationale expander", GREY_DARK, True),
]

STEP_SPACING7 = 0.62
for step_idx, (fi, ti, lbl, col, ret) in enumerate(STEPS_7):
    sy = STEP_START_Y + step_idx * STEP_SPACING7
    if sy > LINE_BOT - 0.05:
        break
    if fi == ti:
        cx = col_cx7[fi]
        rect(s7, cx - 0.6, sy - 0.14, 1.22, 0.38,
             fill=IE_LIGHT, border=col, border_pt=0.8,
             text=lbl, font_size=6.0, text_color=col)
    else:
        seq_arrow7(s7, fi, ti, sy, lbl, col=col, ret=ret)
    txtbox(s7, str(step_idx + 1), 0.05, sy - 0.18, 0.2, 0.25,
           font_size=7, bold=True, color=GREY_MED)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — Data Flow: EU Hub Inbound Logging
# ═══════════════════════════════════════════════════════════════
s8 = new_slide(prs)
slide_title(s8, "Data Flow 3 — EU Hub Inbound Request Logging",
            "Every request from any member state is intercepted, timed and logged")

PARTICIPANTS_8 = [
    ("Member State\nApp (e.g. IE)", EU_MED,  WHITE),
    ("ApiLogging\nMiddleware",       EU_DARK, WHITE),
    ("FastAPI\nRoute Handler",       EU_MED,  WHITE),
    ("eu_vat.db\napi_log",          DB_DARK, WHITE),
    ("EU Hub\nDashboard :8502",     EU_DARK, WHITE),
    ("👤 EU\nAdministrator",        GREY_DARK, WHITE),
]

N8 = len(PARTICIPANTS_8)
COL_W8 = USABLE_W / N8
col_cx8 = [LEFT_MARGIN + (i + 0.5) * COL_W8 for i in range(N8)]

for i, (name, fill, tc) in enumerate(PARTICIPANTS_8):
    rect(s8, LEFT_MARGIN + i * COL_W8 + 0.08, PART_Y,
         COL_W8 - 0.16, PART_H,
         fill=fill, border=fill,
         text=name, font_size=9, bold=True, text_color=tc)
    line(s8, col_cx8[i], LINE_TOP, col_cx8[i], LINE_BOT,
         color=GREY_MED, width_pt=0.75, arrow_end=False, dashed=True)

def seq_arrow8(slide, fi, ti, sy, lbl, col=GREY_DARK, ret=False):
    x1, x2 = col_cx8[fi], col_cx8[ti]
    line(slide, x1, sy, x2, sy,
         color=col, width_pt=1.5,
         arrow_end=(not ret), arrow_start=ret)
    mid_x = (min(x1, x2) + abs(x2 - x1) / 2) - 0.7
    txtbox(slide, lbl, mid_x, sy - 0.3, 1.5, 0.3,
           font_size=8, color=col, align=PP_ALIGN.CENTER)

STEPS_8 = [
    (0, 1, "HTTP request\n+ X-Client-Country: IE",  EU_MED,  False),
    (1, 1, "record t₀",                             EU_DARK, False),
    (1, 2, "forward request",                       EU_DARK, False),
    (2, 3, "query invoices / stats",                DB_MED,  False),
    (3, 2, "results",                               DB_MED,  True),
    (2, 1, "response\n+ X-Records-Returned header", EU_DARK, True),
    (1, 1, "elapsed = now − t₀",                   EU_DARK, False),
    (1, 3, "INSERT api_log\n(timestamp UTC,\nmethod, endpoint,\nclient_country,\nstatus, latency, records)", DB_MED, False),
    (1, 0, "response (pass-through)",               EU_MED,  True),
    (5, 4, "Open Activity Log page",                GREY_DARK, False),
    (4, 3, "get_api_logs(limit=N)",                 DB_MED,  False),
    (3, 4, "log rows",                              DB_MED,  True),
    (4, 5, "Table: timestamp · client\nstatus · latency · records",  EU_DARK, True),
]

STEP_SPACING8 = 0.72
for step_idx, (fi, ti, lbl, col, ret) in enumerate(STEPS_8):
    sy = STEP_START_Y + step_idx * STEP_SPACING8
    if sy > LINE_BOT - 0.05:
        break
    if fi == ti:
        cx = col_cx8[fi]
        rect(s8, cx - 0.65, sy - 0.14, 1.3, 0.38,
             fill=EU_LIGHT, border=col, border_pt=0.8,
             text=lbl, font_size=7, text_color=col)
    else:
        seq_arrow8(s8, fi, ti, sy, lbl, col=col, ret=ret)
    txtbox(s8, str(step_idx + 1), 0.05, sy - 0.18, 0.2, 0.25,
           font_size=8, bold=True, color=GREY_MED)


# ── Save ─────────────────────────────────────────────────────────────────────
out_path = "EU_VAT_Audit_System_Architecture.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
